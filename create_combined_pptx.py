"""Generate Combined Deck: Overall > VIP Analysis > KPI > VIP Solution."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Color palette
BG_DARK = RGBColor(0x1E, 0x27, 0x61)
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF7)
ACCENT = RGBColor(0xE6, 0x3E, 0x31)
ACCENT2 = RGBColor(0x2E, 0xA0, 0x6A)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x1E, 0x1E, 0x2E)
TEXT_MUTED = RGBColor(0x6B, 0x72, 0x80)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_HEADER = RGBColor(0x1E, 0x27, 0x61)
TABLE_ALT = RGBColor(0xEE, 0xEF, 0xF5)
HIGHLIGHT_YELLOW = RGBColor(0xFF, 0xD7, 0x00)
PANEL_BG = RGBColor(0x28, 0x33, 0x78)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
KPI_RED = RGBColor(0xE6, 0x3E, 0x31)
METRIC_GREEN = RGBColor(0x2E, 0xA0, 0x6A)
MECHANIC_BG = RGBColor(0x3A, 0x45, 0x8A)
CAUSE_BG = RGBColor(0xE8, 0xE8, 0xE8)
SOLUTION_GREEN = RGBColor(0x27, 0xAE, 0x60)
LOSS_BG = RGBColor(0xFC, 0xE4, 0xEC)


def add_shape(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_rounded_box(slide, left, top, width, height, fill_color, border_color=None, border=None):
    shape = slide.shapes.add_shape(5, left, top, width, height)  # 5 = rounded rect
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    bc = border_color or border
    if bc:
        shape.line.color.rgb = bc
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def set_cell(cell, text, font_size=10, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT, fill=None):
    cell.text = str(text)
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill


def add_table(slide, rows_data, headers, left, top, width, col_widths=None, font_size=9):
    n_rows = len(rows_data) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.3 * n_rows))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for i, h in enumerate(headers):
        set_cell(table.cell(0, i), h, font_size=font_size, bold=True,
                 color=TEXT_WHITE, align=PP_ALIGN.CENTER, fill=TABLE_HEADER)
    for r, row in enumerate(rows_data):
        fill = TABLE_ALT if r % 2 == 0 else CARD_BG
        for c, val in enumerate(row):
            align = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            set_cell(table.cell(r + 1, c), val, font_size=font_size,
                     color=TEXT_DARK, align=align, fill=fill)
    return table_shape


def add_text(slide, text, left, top, width, height, font_size=14, bold=False,
             color=TEXT_DARK, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_multiline(slide, lines, left, top, width, height, font_size=12,
                  color=TEXT_DARK, line_spacing=1.3, bold_first=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(font_size * (line_spacing - 1))
        if bold_first and i == 0:
            p.font.bold = True
    return txBox


def title_bar(slide, title_text=None):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_LIGHT
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), BG_DARK)
    if title_text:
        add_text(slide, title_text, Inches(0.6), Inches(0.2), Inches(12), Inches(0.6),
                 font_size=26, bold=True, color=TEXT_WHITE)


def bullet_card(slide, x, y, num, title, desc, col):
    add_text(slide, num, x, y, Inches(0.5), Inches(0.4), font_size=20, bold=True, color=col)
    add_text(slide, title, x + Inches(0.5), y, Inches(4.3), Inches(0.35), font_size=16, bold=True, color=TEXT_WHITE)
    add_text(slide, desc, x + Inches(0.5), y + Inches(0.35), Inches(4.3), Inches(0.3),
             font_size=11, color=RGBColor(0xAA, 0xB0, 0xCC))


def add_line_h(slide, x, y, length, color=TEXT_MUTED):
    add_shape(slide, x, y, length, Inches(0.025), color)


def add_line_v(slide, x, y, length, color=TEXT_MUTED):
    add_shape(slide, x, y, Inches(0.025), length, color)


def add_arrow_right(slide, x, y, color=TEXT_MUTED):
    """Small right-pointing triangle."""
    add_text(slide, "\u25B6", x, y - Inches(0.08), Inches(0.2), Inches(0.2),
             font_size=8, color=color, align=PP_ALIGN.CENTER)


def kpi_box(slide, x, y, w, h, label, current, target, bg=CARD_BG, border=None):
    """Draw a KPI metric box with label, current value, and KPI target."""
    add_rounded_box(slide, x, y, w, h, bg, border)
    add_text(slide, label, x + Inches(0.08), y + Inches(0.02), w - Inches(0.16), Inches(0.22),
             font_size=9, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, current, x + Inches(0.08), y + Inches(0.22), w - Inches(0.16), Inches(0.22),
             font_size=12, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, target, x + Inches(0.08), y + Inches(0.45), w - Inches(0.16), Inches(0.20),
             font_size=9, bold=True, color=KPI_RED, align=PP_ALIGN.CENTER)


def mechanic_box(slide, x, y, w, h, title, bg=MECHANIC_BG):
    add_rounded_box(slide, x, y, w, h, bg)
    add_text(slide, title, x + Inches(0.08), y + Inches(0.05), w - Inches(0.16), h - Inches(0.1),
             font_size=10, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)


def cause_box(slide, x, y, w, h, text, bg=CAUSE_BG):
    add_rounded_box(slide, x, y, w, h, bg)
    add_text(slide, text, x + Inches(0.06), y + Inches(0.04), w - Inches(0.12), h - Inches(0.08),
             font_size=8, color=TEXT_DARK, align=PP_ALIGN.LEFT)


def solution_box(slide, x, y, w, h, text, bg=SOLUTION_GREEN):
    add_rounded_box(slide, x, y, w, h, bg)
    add_text(slide, text, x + Inches(0.06), y + Inches(0.04), w - Inches(0.12), h - Inches(0.08),
             font_size=9, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ==============================================================================
# PART 1: OVERALL ANALYSIS (slides 1-4)
# ==============================================================================

# SLIDE 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = BG_DARK
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT)
add_text(slide, "PHÂN TÍCH MONETIZATION & VIP UPDATE", Inches(1), Inches(2.0), Inches(11), Inches(1),
         font_size=42, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER, font_name="Arial Black")
add_text(slide, "Cờ Tỷ Phú (CTP)", Inches(1), Inches(2.8), Inches(11), Inches(0.7),
         font_size=28, color=HIGHLIGHT_YELLOW, align=PP_ALIGN.CENTER)
add_text(slide, "Overall Analysis → VIP Deep-dive → KPI Target → Solution Design",
         Inches(1), Inches(3.8), Inches(11), Inches(0.5),
         font_size=16, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "Data: 01/02/2026 - 01/04/2026", Inches(1), Inches(4.5), Inches(11), Inches(0.5),
         font_size=14, color=RGBColor(0xAA, 0xB0, 0xCC), align=PP_ALIGN.CENTER)
add_shape(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), ACCENT)

# SLIDE 2: Overall Revenue Breakdown
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "Tổng Quan Monetization — T3/2026")
add_table(slide, [
    ["Tổng", "1.52 tỷ", "100%", "4,603", "330k"],
    ["GEM", "875M", "57.5%", "1,173", "746k"],
    ["Gold (gộp)", "379M", "24.9%", "3,199", "119k"],
    ["First Pay", "131M", "8.6%", "1,570", "83k"],
    ["VIP", "87M", "5.8%", "724", "121k"],
    ["SEASON_PASS", "46M", "3.0%", "387", "118k"],
], ["Sản phẩm", "Revenue", "Share", "PU", "ARPPU"],
    Inches(0.5), Inches(1.3), Inches(7.0), col_widths=[2.0, 1.2, 1.0, 1.0, 1.0], font_size=11)

add_table(slide, [
    ["Tổng", "1.16 tỷ", "1.52 tỷ", "+31%"],
    ["GEM", "888M", "875M", "-1.5%"],
    ["Gold (gộp)", "97M", "379M", "+291%"],
    ["First Pay", "72M", "131M", "+82%"],
    ["VIP", "55M", "87M", "+58%"],
    ["SEASON_PASS", "29M", "46M", "+58%"],
], ["Sản phẩm", "Rev T2", "Rev T3", "Tăng trưởng"],
    Inches(0.5), Inches(4.5), Inches(7.0), col_widths=[2.0, 1.5, 1.5, 1.5], font_size=11)

add_shape(slide, Inches(8.0), Inches(1.3), Inches(4.8), Inches(5.5), CARD_BG)
for i, (t, d, c) in enumerate([
    ("GEM giảm 1.5%", "Dù PU tăng 16%, ARPPU giảm.\nNguồn rev chính đang yếu đi.", ACCENT),
    ("Gold/First Pay +291%", "Tăng trưởng hoàn toàn từ\nnew user do UA investment.", ACCENT2),
    ("VIP +58%", "Tăng ngang UA growth, không\noutperform. Room to improve lớn nhất.", BG_DARK),
]):
    y = Inches(1.6 + i * 1.7)
    add_text(slide, "●", Inches(8.3), y, Inches(0.3), Inches(0.3), font_size=16, bold=True, color=c)
    add_text(slide, t, Inches(8.7), y, Inches(3.8), Inches(0.35), font_size=14, bold=True, color=TEXT_DARK)
    add_text(slide, d, Inches(8.7), y + Inches(0.4), Inches(3.8), Inches(0.7), font_size=11, color=TEXT_MUTED)

# SLIDE 3: Repay Rate 30 ngày
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "Repay Rate 30 Ngày — Vấn Đề Lớn Của Game")
add_shape(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(2.2), CARD_BG)
add_text(slide, "72.6%", Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.8),
         font_size=60, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, "payer mất sau 30 ngày", Inches(0.5), Inches(2.3), Inches(5.5), Inches(0.4),
         font_size=18, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
add_text(slide, "Chỉ 625 / 2,282 user T2 quay lại nạp trong 30 ngày sau lần nạp cuối",
         Inches(0.5), Inches(2.8), Inches(5.5), Inches(0.4), font_size=12, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
add_table(slide, [
    ["Tổng", "2,282", "625 (27.4%)", "1,657 (72.6%)"],
    ["<7d (New)", "640", "141 (22.0%)", "499 (78.0%)"],
    ["7-30d", "430", "94 (21.9%)", "336 (78.1%)"],
    ["30-90d", "376", "98 (26.1%)", "278 (73.9%)"],
    ["90-180d", "245", "75 (30.6%)", "170 (69.4%)"],
    [">180d (Core)", "591", "217 (36.7%)", "374 (63.3%)"],
], ["Segment", "Users", "Repay", "Mất"],
    Inches(7), Inches(1.3), Inches(5.8), col_widths=[1.8, 0.8, 1.5, 1.5], font_size=11)

add_shape(slide, Inches(0.5), Inches(4.0), Inches(12.3), Inches(2.8), CARD_BG)
add_multiline(slide, [
    "Đo từ lần nạp cuối trong T2, trong 30 ngày sau có nạp lại không?",
    "",
    "• New user mất 78%, Core vẫn mất 63% — cả game đều có vấn đề",
    "• Game đang sống bằng UA — liên tục đổ user mới vào thay thế user cũ bỏ",
    "• Cần cơ chế recurring để giữ payer lâu hơn",
], Inches(1.0), Inches(4.2), Inches(11.3), Inches(2.2),
              font_size=14, color=TEXT_DARK, line_spacing=1.4)

# SLIDE 3b: 52% chỉ nạp 1 lần trong tháng
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "52% Payer Chỉ Nạp 1 Lần Trong Tháng")
add_shape(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(2.2), CARD_BG)
add_text(slide, "52%", Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.8),
         font_size=60, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, "payer chỉ nạp đúng 1 lần trong T2", Inches(0.5), Inches(2.3), Inches(5.5), Inches(0.4),
         font_size=18, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
add_text(slide, "1,187 / 2,282 payer T2 không nạp thêm lần nào",
         Inches(0.5), Inches(2.8), Inches(5.5), Inches(0.4), font_size=12, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

add_table(slide, [
    ["1 lần duy nhất", "1,187", "52.0%"],
    ["2+ lần", "1,095", "48.0%"],
    ["Tổng", "2,282", "100%"],
], ["Số lần nạp trong T2", "Users", "%"],
    Inches(7), Inches(1.3), Inches(5.5), col_widths=[2.2, 1.5, 1.0], font_size=11)

add_shape(slide, Inches(0.5), Inches(4.0), Inches(12.3), Inches(2.8), CARD_BG)
add_multiline(slide, [
    "Hơn nửa payer nạp 1 lần rồi không quay lại trong cả tháng",
    "",
    "• Kết hợp với repay rate 27.4% → game không convert được payer thành repeat",
    "• Revenue phụ thuộc vào liên tục acquire payer mới, không phải giữ payer cũ",
    "• VIP subscription là cơ hội tạo habit nạp lại hàng tháng",
], Inches(1.0), Inches(4.2), Inches(11.3), Inches(2.2),
              font_size=14, color=TEXT_DARK, line_spacing=1.4)

# SLIDE 4: VIP Opportunity (2 lý do)
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "Tại Sao Update VIP Là Cơ Hội Lớn?")
for i, (t, d, c, n) in enumerate([
    ("VIP kém nhất\ntrong các mặt hàng",
     "VIP +58% vs Gold +291%.\nPerformance thấp nhất = room\nto improve lớn nhất.\nCải thiện cái kém dễ hơn\ncải thiện cái đang tốt.\n\n"
     "VIP share chỉ 5.75% total rev\n→ Potential tăng rất lớn nếu\nredesign đúng hướng.", ACCENT, "1"),
    ("Repay rate game\nchỉ 27.4%\n+ 52% nạp 1 lần",
     "72.6% payer mất sau 30 ngày.\n52% chỉ nạp 1 lần trong tháng.\n\nVIP subscription là cơ chế\nrecurring hiếm hoi của game.\nRedesign VIP → tạo habit repay\n→ kéo repay rate tổng lên\n→ cải thiện performance game.", BG_DARK, "2"),
]):
    x = Inches(0.5 + i * 6.2)
    add_shape(slide, x, Inches(1.3), Inches(5.8), Inches(5.5), CARD_BG)
    add_text(slide, n, x + Inches(0.2), Inches(1.5), Inches(0.6), Inches(0.6),
             font_size=28, bold=True, color=c, align=PP_ALIGN.CENTER)
    add_text(slide, t, x + Inches(0.9), Inches(1.5), Inches(4.5), Inches(1.0),
             font_size=18, bold=True, color=TEXT_DARK)
    add_text(slide, d, x + Inches(0.4), Inches(2.8), Inches(5.0), Inches(3.5),
             font_size=14, color=TEXT_MUTED)

# Arrow connecting 1+2 → VIP analysis
add_shape(slide, Inches(4.5), Inches(6.0), Inches(4.3), Inches(0.6), ACCENT2)
add_text(slide, "→ Phân tích & Update VIP là cơ hội tăng performance game",
         Inches(4.5), Inches(6.05), Inches(4.3), Inches(0.5),
         font_size=13, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

# ==============================================================================
# PART 2: VIP ANALYSIS (slides 5-10)
# ==============================================================================

# SLIDE 5: VIP Revenue
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "Tổng Quan Doanh Thu VIP")
for i, (num, label, col) in enumerate([
    ("4.78%", "T2 VIP Share", TEXT_MUTED), ("5.75%", "T3 VIP Share", ACCENT2),
    ("+63%", "Tăng trưởng PU", ACCENT2), ("55.3M", "Doanh thu T2", TEXT_MUTED),
    ("87.4M", "Doanh thu T3", ACCENT),
]):
    x = Inches(0.5 + i * 2.5)
    add_shape(slide, x, Inches(1.4), Inches(2.2), Inches(1.3), CARD_BG)
    add_text(slide, num, x, Inches(1.55), Inches(2.2), Inches(0.6),
             font_size=32, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, Inches(2.1), Inches(2.2), Inches(0.3),
             font_size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
add_table(slide, [
    ["T2/2026", "1.16 tỷ", "55.3M", "4.78%", "445", "2,282"],
    ["T3/2026", "1.52 tỷ", "87.4M", "5.75%", "724", "4,603"],
], ["Tháng", "Tổng Rev", "VIP Rev", "VIP %", "VIP PU", "Tổng PU"],
    Inches(0.5), Inches(3.2), Inches(12.3), col_widths=[1.8, 2.0, 2.0, 1.5, 1.5, 1.5], font_size=11)
add_text(slide, "VIP PU tăng từ 445 lên 724, nhưng do UA investment tăng → A1/N1 tăng.\n"
         "Chưa thể kết luận VIP tự cải thiện, cần normalize theo A1/N1.",
         Inches(0.5), Inches(4.8), Inches(12), Inches(0.7), font_size=13, color=TEXT_DARK)

# SLIDE 6: VIP Buyers by Age
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "Người Mua VIP Theo Tuổi Tài Khoản")
add_table(slide, [
    ["<7 ngày (Mới)", "361", "33.1%", "42.8M", "118k", "1.17"],
    ["7-30 ngày (Sớm)", "207", "19.0%", "24.5M", "118k", "1.17"],
    ["30-90 ngày (TB)", "172", "15.8%", "21.5M", "125k", "1.19"],
    ["90-180 ngày (Trung thành)", "100", "9.2%", "14.4M", "144k", "1.43"],
    [">180 ngày (Core)", "252", "23.1%", "44.4M", "176k", "1.74"],
], ["Phân khúc", "Số người", "%", "Doanh thu", "ARPPU", "TB lần mua"],
    Inches(0.5), Inches(1.4), Inches(12.3), col_widths=[3.0, 1.5, 1.0, 2.0, 1.5, 1.8], font_size=11)
for i, (num, desc, col) in enumerate([
    ("33.1%", "User mới chiếm PU nhiều nhất\nnhưng gần như không renewal (1.17 lần)", ACCENT),
    ("23.1%", "Core >180d chỉ 23% PU nhưng\n30% rev, repeat cao nhất (1.74x)", ACCENT2),
    ("176k", "ARPPU tăng dần theo tuổi:\n118k → 125k → 144k → 176k", BG_DARK),
]):
    x = Inches(0.5 + i * 4.2)
    add_shape(slide, x, Inches(4.4), Inches(3.8), Inches(2.4), CARD_BG)
    add_text(slide, num, x, Inches(4.55), Inches(3.8), Inches(0.6),
             font_size=36, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(slide, desc, x + Inches(0.3), Inches(5.2), Inches(3.2), Inches(1.2),
             font_size=11, color=TEXT_DARK, align=PP_ALIGN.CENTER)

# SLIDE 7: VIP Renewal
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "82.8% Chỉ Mua VIP 1 Lần — Vấn Đề Lớn Nhất")
add_shape(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(2.5), CARD_BG)
add_text(slide, "82.8%", Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.9),
         font_size=60, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, "user chỉ mua VIP 1 lần rồi bỏ", Inches(0.5), Inches(2.4), Inches(5.5), Inches(0.5),
         font_size=18, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
add_text(slide, "Chỉ 17.2% mua lại lần 2+", Inches(1.0), Inches(3.0), Inches(4.5), Inches(0.4),
         font_size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
add_table(slide, [
    ["1 lần", "904", "82.8%"], ["2 lần", "106", "9.7%"],
    ["3+ lần", "82", "7.5%"],
], ["Số lần mua", "Số người", "%"],
    Inches(7), Inches(1.3), Inches(5.5), col_widths=[2.0, 1.5, 1.5], font_size=11)

# Deep-dive funnel
add_shape(slide, Inches(0.5), Inches(4.2), Inches(7.0), Inches(3.0), CARD_BG)
add_text(slide, "Phễu sau khi VIP hết hạn", Inches(0.8), Inches(4.3), Inches(6.4), Inches(0.3),
         font_size=14, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
for left, top, width, label, color in [
    (Inches(0.8), Inches(4.8), Inches(6.4), "1,092 VIP Buyers (100%)", BG_DARK),
    (Inches(1.2), Inches(5.35), Inches(5.6), "904 mua 1 lần (82.8%)", ACCENT),
    (Inches(1.6), Inches(5.9), Inches(4.8), "~700 vẫn chơi game (77%)", RGBColor(0xE6, 0x7E, 0x22)),
    (Inches(2.0), Inches(6.45), Inches(4.0), "~593 không nạp gì (85%)", RGBColor(0xC0, 0x39, 0x2B)),
]:
    add_shape(slide, left, top, width, Inches(0.45), color)
    add_text(slide, label, left + Inches(0.1), top + Inches(0.05), width - Inches(0.2), Inches(0.35),
             font_size=11, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

add_shape(slide, Inches(8.0), Inches(4.2), Inches(4.8), Inches(3.0), CARD_BG)
add_text(slide, "Phát hiện chính", Inches(8.3), Inches(4.3), Inches(4.2), Inches(0.3),
         font_size=14, bold=True, color=BG_DARK)
for i, (t, d, c) in enumerate([
    ("77% vẫn chơi nhưng 85% ngừng nạp", "VIP = cửa cuối trước F2P", ACCENT),
    ("Window kích repay = 15 ngày", "Sau 15 ngày gần như mất luôn", ACCENT2),
    ("593 user = pool win-back lớn nhất", "Đã từng chi tiền, barrier thấp", BG_DARK),
]):
    y = Inches(4.8 + i * 0.8)
    add_text(slide, "●", Inches(8.3), y, Inches(0.3), Inches(0.2), font_size=12, bold=True, color=c)
    add_text(slide, t, Inches(8.6), y, Inches(3.8), Inches(0.25), font_size=11, bold=True, color=TEXT_DARK)
    add_text(slide, d, Inches(8.6), y + Inches(0.25), Inches(3.8), Inches(0.25), font_size=10, color=TEXT_MUTED)

# ==============================================================
# SLIDE 8b: Chốt vấn đề VIP
# ==============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "Chốt Vấn Đề VIP — 2 KPI Cần Fix")

add_text(slide, "Từ phân tích trên, VIP hiện tại có 2 vấn đề cốt lõi:",
         Inches(0.6), Inches(1.1), Inches(12), Inches(0.4),
         font_size=16, bold=True, color=TEXT_DARK)

# Vấn đề 1: Repay trong tháng
add_shape(slide, Inches(0.5), Inches(1.8), Inches(6.0), Inches(4.8), CARD_BG)
add_text(slide, "1", Inches(0.7), Inches(1.9), Inches(0.5), Inches(0.5),
         font_size=28, bold=True, color=ACCENT)
add_text(slide, "Không repay trong tháng", Inches(1.2), Inches(1.95), Inches(5.0), Inches(0.4),
         font_size=18, bold=True, color=TEXT_DARK)
add_text(slide, "ARPT hiện tại: 1.0 — user mua VIP 1 lần rồi không mua tiếp trong tháng",
         Inches(0.8), Inches(2.5), Inches(5.4), Inches(0.4),
         font_size=12, bold=True, color=ACCENT)

add_multiline(slide, [
    "Root causes:",
    "• VIP vs Free như nhau về hình ảnh — không có identity",
    "• x2 EXP chỉ cần khi upgrade, xong thì hết nhu cầu",
    "• Mua lần 2 không có gì hơn lần 1 — không có progression",
    "• Không có đặc quyền exclusive cho VIP",
    "",
    "→ Không có lý do mua lại trong cùng tháng",
], Inches(0.8), Inches(3.1), Inches(5.4), Inches(3.0),
              font_size=12, color=TEXT_DARK, line_spacing=1.4)

# Vấn đề 2: Repay tháng sau
add_shape(slide, Inches(6.8), Inches(1.8), Inches(6.0), Inches(4.8), CARD_BG)
add_text(slide, "2", Inches(7.0), Inches(1.9), Inches(0.5), Inches(0.5),
         font_size=28, bold=True, color=BG_DARK)
add_text(slide, "Không repay tháng sau", Inches(7.5), Inches(1.95), Inches(5.0), Inches(0.4),
         font_size=18, bold=True, color=TEXT_DARK)
add_text(slide, "Repeat rate hiện tại: 17% — 82.8% mua 1 lần rồi bỏ hẳn",
         Inches(7.1), Inches(2.5), Inches(5.4), Inches(0.4),
         font_size=12, bold=True, color=BG_DARK)

add_multiline(slide, [
    "Root causes:",
    "• 2 lần mua độc lập, không có streak hay loyalty",
    "• Không có noti khi gần hết VIP",
    "• Tháng này mua VIP không khác gì tháng trước",
    "  — không có new content hay progression cross-tháng",
    "",
    "→ Không có lý do quay lại mua tháng tiếp theo",
], Inches(7.1), Inches(3.1), Inches(5.4), Inches(3.0),
              font_size=12, color=TEXT_DARK, line_spacing=1.4)

# ==============================================================================
# PART 3: KPI TARGET & SOLUTION
# ==============================================================================

# Section Divider
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = BG_DARK
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT2)
add_text(slide, "VIP UPDATE", Inches(1), Inches(2.5), Inches(11), Inches(0.8),
         font_size=44, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER, font_name="Arial Black")
add_text(slide, "KPI Target → Solution Design", Inches(1), Inches(3.3), Inches(11), Inches(0.5),
         font_size=22, color=HIGHLIGHT_YELLOW, align=PP_ALIGN.CENTER)
add_shape(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), ACCENT2)

# ==============================================================
# SLIDE 9: KPI Breakdown Tree (page 1 - top level)
# ==============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "KPI Framework — Phân rã target Rev 350tr/tháng")

# --- LEVEL 0: Rev ---
kpi_box(slide, Inches(0.3), Inches(2.8), Inches(1.4), Inches(0.7),
        "Rev/tháng", "~72tr", "KPI: 350tr", bg=RGBColor(0xE8, 0xFF, 0xE8), border=ACCENT2)

# Line from Rev → split
add_line_h(slide, Inches(1.7), Inches(3.15), Inches(0.5))

# --- LEVEL 1: PU + ARPPU ---
# PU
kpi_box(slide, Inches(2.2), Inches(1.6), Inches(1.3), Inches(0.7),
        "PU", "724", "KPI: 1,162", bg=CARD_BG, border=BG_DARK)
# ARPPU
kpi_box(slide, Inches(2.2), Inches(3.9), Inches(1.3), Inches(0.7),
        "ARPPU", "100k", "KPI: 300k", bg=CARD_BG, border=BG_DARK)

# Vertical line connecting PU and ARPPU
add_line_v(slide, Inches(2.18), Inches(2.0), Inches(2.3))
add_line_h(slide, Inches(2.18), Inches(2.0), Inches(0.05))
add_line_h(slide, Inches(2.18), Inches(4.25), Inches(0.05))

# --- LEVEL 2 from PU ---
add_line_h(slide, Inches(3.5), Inches(1.95), Inches(0.4))

# New PU
kpi_box(slide, Inches(3.9), Inches(1.15), Inches(1.6), Inches(0.7),
        "New PU: 82%", "594", "KPI: 800 (+200)", bg=CARD_BG, border=ACCENT2)
# Repay next month
kpi_box(slide, Inches(3.9), Inches(2.3), Inches(1.6), Inches(0.7),
        "Repay NM: 17%", "130", "KPI: 362 (+232)", bg=CARD_BG, border=ACCENT)

add_line_v(slide, Inches(3.88), Inches(1.5), Inches(1.3))
add_line_h(slide, Inches(3.88), Inches(1.5), Inches(0.05))
add_line_h(slide, Inches(3.88), Inches(2.65), Inches(0.05))

# --- LEVEL 2 from ARPPU ---
add_line_h(slide, Inches(3.5), Inches(4.25), Inches(0.4))

# ARPT
kpi_box(slide, Inches(3.9), Inches(3.9), Inches(1.3), Inches(0.7),
        "ARPT", "1.0", "KPI: 3", bg=CARD_BG, border=ACCENT)
# ARPPT
kpi_box(slide, Inches(3.9), Inches(4.8), Inches(1.3), Inches(0.7),
        "ARPPT", "100k", "Giữ nguyên", bg=CARD_BG, border=TEXT_MUTED)

add_line_v(slide, Inches(3.88), Inches(4.25), Inches(0.95))
add_line_h(slide, Inches(3.88), Inches(4.25), Inches(0.05))
add_line_h(slide, Inches(3.88), Inches(5.15), Inches(0.05))

# --- LEVEL 3: Mechanic boxes ---
# From New PU → Mechanic convert
add_line_h(slide, Inches(5.5), Inches(1.5), Inches(0.4))
mechanic_box(slide, Inches(5.9), Inches(1.15), Inches(1.8), Inches(0.7),
             "Mechanic\nConvert New Pay")

# From Repay NM → Mechanic repay next month
add_line_h(slide, Inches(5.5), Inches(2.65), Inches(0.4))
mechanic_box(slide, Inches(5.9), Inches(2.3), Inches(1.8), Inches(0.7),
             "Mechanic\nRepay Next Month")

# From ARPT → Mechanic repay trong tháng
add_line_h(slide, Inches(5.2), Inches(4.25), Inches(0.7))
mechanic_box(slide, Inches(5.9), Inches(3.9), Inches(1.8), Inches(0.7),
             "Mechanic\nRepay Trong Tháng")

# --- LEVEL 4: Root causes → Solutions ---
# Convert New Pay causes
add_line_h(slide, Inches(7.7), Inches(1.3), Inches(0.3))
cause_box(slide, Inches(8.0), Inches(1.1), Inches(2.3), Inches(0.4),
          "58% (223) chưa mua VIP")
solution_box(slide, Inches(10.5), Inches(1.1), Inches(2.5), Inches(0.4),
             "Offer lần đầu")
add_line_h(slide, Inches(10.3), Inches(1.3), Inches(0.2))

cause_box(slide, Inches(8.0), Inches(1.6), Inches(2.3), Inches(0.4),
          "42% (159) lapsed"  )
solution_box(slide, Inches(10.5), Inches(1.6), Inches(2.5), Inches(0.4),
             "Offer welcome back VIP")
add_line_h(slide, Inches(10.3), Inches(1.8), Inches(0.2))

add_line_v(slide, Inches(7.98), Inches(1.3), Inches(0.5))

# Repay Next Month causes
add_line_h(slide, Inches(7.7), Inches(2.5), Inches(0.3))
cause_box(slide, Inches(8.0), Inches(2.15), Inches(2.3), Inches(0.38),
          "Mua độc lập, ko streak")
solution_box(slide, Inches(10.5), Inches(2.15), Inches(2.5), Inches(0.38),
             "Streak tháng")
add_line_h(slide, Inches(10.3), Inches(2.34), Inches(0.2))

cause_box(slide, Inches(8.0), Inches(2.58), Inches(2.3), Inches(0.38),
          "Ko noti khi gần hết VIP")
solution_box(slide, Inches(10.5), Inches(2.58), Inches(2.5), Inches(0.38),
             "Noti những ngày cuối")
add_line_h(slide, Inches(10.3), Inches(2.77), Inches(0.2))

cause_box(slide, Inches(8.0), Inches(3.01), Inches(2.3), Inches(0.38),
          "Ko new content mỗi tháng")
solution_box(slide, Inches(10.5), Inches(3.01), Inches(2.5), Inches(0.38),
             "Renew content")
add_line_h(slide, Inches(10.3), Inches(3.2), Inches(0.2))

add_line_v(slide, Inches(7.98), Inches(2.34), Inches(0.86))

# Repay Trong Tháng causes
add_line_h(slide, Inches(7.7), Inches(4.1), Inches(0.3))

cause_box(slide, Inches(8.0), Inches(3.7), Inches(2.3), Inches(0.38),
          "VIP vs Free như nhau hình ảnh")
solution_box(slide, Inches(10.5), Inches(3.7), Inches(2.5), Inches(0.38),
             "Cosmetic")
add_line_h(slide, Inches(10.3), Inches(3.89), Inches(0.2))

cause_box(slide, Inches(8.0), Inches(4.13), Inches(2.3), Inches(0.38),
          "x2 EXP chỉ cần khi upgrade")
solution_box(slide, Inches(10.5), Inches(4.13), Inches(2.5), Inches(0.38),
             "Benefit tăng dần")
add_line_h(slide, Inches(10.3), Inches(4.32), Inches(0.2))

cause_box(slide, Inches(8.0), Inches(4.56), Inches(2.3), Inches(0.38),
          "Mua lần 2 = lần 1, ko progression")
solution_box(slide, Inches(10.5), Inches(4.56), Inches(2.5), Inches(0.38),
             "Streak benefit")
add_line_h(slide, Inches(10.3), Inches(4.75), Inches(0.2))

cause_box(slide, Inches(8.0), Inches(4.99), Inches(2.3), Inches(0.38),
          "Ko có đặc quyền exclusive")
solution_box(slide, Inches(10.5), Inches(4.99), Inches(2.5), Inches(0.38),
             "Limited access")
add_line_h(slide, Inches(10.3), Inches(5.18), Inches(0.2))

add_line_v(slide, Inches(7.98), Inches(3.89), Inches(1.29))

# Loss aversion label
add_rounded_box(slide, Inches(10.5), Inches(5.6), Inches(2.5), Inches(0.5),
                LOSS_BG, border=ACCENT)
add_text(slide, "Benefit cảm giác mất mát", Inches(10.5), Inches(5.65), Inches(2.5), Inches(0.4),
         font_size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
# Arrows from cosmetic/streak/limited to loss box
add_line_v(slide, Inches(11.75), Inches(5.37), Inches(0.23))

# Legend
add_text(slide, "Xanh = KPI target    Xám = Root cause    Xanh lá = Solution",
         Inches(0.3), Inches(6.2), Inches(6), Inches(0.3),
         font_size=9, color=TEXT_MUTED)

# ==============================================================
# SLIDE 10: Impact Estimation (updated model)
# ==============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "KPI Impact Estimation")

add_text(slide, "Base: T3 có 724 VIP PU (594 new + 130 returning). Target New PU: 800/tháng.",
         Inches(0.6), Inches(1.1), Inches(12), Inches(0.3),
         font_size=14, bold=True, color=TEXT_DARK)
add_text(slide, "Model: New retain 50%, Old retain 80%, New ARPT=2, Old ARPT=3, ARPPT=100k",
         Inches(0.6), Inches(1.45), Inches(12), Inches(0.3),
         font_size=12, color=TEXT_MUTED)

add_table(slide, [
    ["T1", "800", "362 (50%×724)", "1,162", "1,600+1,086 = 2,686", "269M"],
    ["T2", "800", "690 (400+80%×362)", "1,490", "1,600+2,070 = 3,670", "367M"],
    ["T3", "800", "952 (400+80%×690)", "1,752", "1,600+2,856 = 4,456", "446M"],
    ["T4", "800", "1,162 (400+80%×952)", "1,962", "1,600+3,486 = 5,086", "509M"],
    ["T6+", "800", "~2,000", "~2,800", "1,600+6,000 = ~7,600", "~760M"],
], ["Tháng", "PU mới", "Old retained", "Active PU", "Lượt (new×2 + old×3)", "Rev"],
    Inches(0.3), Inches(1.9), Inches(12.7),
    col_widths=[0.8, 0.8, 1.8, 1.0, 2.5, 0.8], font_size=10)

add_shape(slide, Inches(7.5), Inches(4.2), Inches(5.3), Inches(1.5), PANEL_BG)
add_text(slide, "Đạt 350M ngay T2", Inches(7.8), Inches(4.3), Inches(4.7), Inches(0.3),
         font_size=16, bold=True, color=HIGHLIGHT_YELLOW)
add_text(slide, "Pool old user tích lũy (80% retain) → compound\nOld pool converge ~2,000 → VIP thành revenue pillar",
         Inches(7.8), Inches(4.7), Inches(4.7), Inches(0.7),
         font_size=12, color=TEXT_WHITE)

add_text(slide, "Key Assumptions:", Inches(0.6), Inches(4.2), Inches(6), Inches(0.3),
         font_size=14, bold=True, color=BG_DARK)

add_table(slide, [
    ["New VIP PU", "800/tháng", "Rương cam hook + offer lần đầu", "T3: 594, +rương cam kích thêm"],
    ["New retain → old", "50%", "VIP Level streak + countdown", "Hiện tại ~17%, target gấp 3x"],
    ["Old retain rate", "80%", "Loss aversion (cosmetic+DKXX+bonus G)", "VIP Level sunk cost giữ chân"],
    ["New ARPT", "2", "Benefit tăng dần kích mua kỳ 2", "VIP 10 ngày, mua 2 = 20 ngày"],
    ["Old ARPT", "3", "Streak + exclusive + countdown", "Streak benefit kích mua cả tháng"],
], ["KPI", "Target", "Mechanic giải quyết", "Cơ sở"],
    Inches(0.3), Inches(4.7), Inches(12.7),
    col_widths=[1.3, 0.9, 3.2, 2.8], font_size=9)

# Risks
add_shape(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(1.0), CARD_BG)
add_text(slide, "Risks:", Inches(0.8), Inches(6.25), Inches(1), Inches(0.3),
         font_size=11, bold=True, color=ACCENT)
add_text(slide, "Bonus %G migration gây phản ứng  •  Countdown 1 ngày quá gắt cho casual  •  DKXX boost ảnh hưởng competitive balance",
         Inches(1.7), Inches(6.25), Inches(10.8), Inches(0.6),
         font_size=10, color=TEXT_DARK)

# ==============================================================
# SLIDE 11: VIP Level System Rules + Loss Aversion
# ==============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "VIP Level System — Rules")

add_text(slide, "Mỗi lần mua VIP liên tiếp = +1 Level. Max Lv5. Không mua = Reset về Lv0.",
         Inches(0.6), Inches(1.1), Inches(12), Inches(0.4),
         font_size=16, bold=True, color=ACCENT)

# Rules
add_shape(slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(5.2), CARD_BG)
add_text(slide, "Rule chung", Inches(0.9), Inches(1.8), Inches(5.2), Inches(0.4),
         font_size=16, bold=True, color=BG_DARK)
add_multiline(slide, [
    "1. VIP có 5 cấp (Lv1-Lv5). Mỗi cấp unlock quyền lợi mới",
    "2. Mỗi VIP trị giá 100k, kéo dài 10 ngày",
    "3. Hết 10 ngày → countdown 1 ngày để nạp level tiếp",
    "4. Không nạp trong countdown → reset về Lv0",
    "5. Level 5: countdown 3 ngày (ưu đãi core)",
    "6. Timeline 10 ngày đi riêng theo user",
    "7. Chỉ thấy benefit level tiếp theo +1",
    "   (Lv1 thấy quà Lv2, Lv3+ bị ẩn)",
], Inches(0.9), Inches(2.3), Inches(5.2), Inches(3.5),
              font_size=12, color=TEXT_DARK, line_spacing=1.5)

# Lapsed
add_shape(slide, Inches(6.8), Inches(1.7), Inches(6.0), Inches(1.5), RGBColor(0x3A, 0x45, 0x8A))
add_text(slide, "Rule cho Lapsed", Inches(7.1), Inches(1.8), Inches(5.4), Inches(0.3),
         font_size=13, bold=True, color=ORANGE)
add_multiline(slide, [
    "Qua 1 mùa không nạp VIP → popup KM skip lên Lv2 luôn",
    "Áp dụng 1 lần/user (anti-gaming)",
], Inches(7.1), Inches(2.2), Inches(5.4), Inches(0.5),
              font_size=11, color=TEXT_WHITE)

# New user
add_shape(slide, Inches(6.8), Inches(3.4), Inches(6.0), Inches(1.5), RGBColor(0x3A, 0x45, 0x8A))
add_text(slide, "Rule cho New User", Inches(7.1), Inches(3.5), Inches(5.4), Inches(0.3),
         font_size=13, bold=True, color=ACCENT2)
add_multiline(slide, [
    "58% payer chưa mua VIP → Offer lần đầu (giá KM hoặc bonus)",
    "Rương cam ở Lv1 = hook mạnh cho new user",
], Inches(7.1), Inches(3.9), Inches(5.4), Inches(0.5),
              font_size=11, color=TEXT_WHITE)

# Loss aversion
add_shape(slide, Inches(6.8), Inches(5.1), Inches(6.0), Inches(1.8), PANEL_BG)
add_text(slide, "Loss Aversion — mất VIP đau ở 3 chiều", Inches(7.1), Inches(5.2), Inches(5.4), Inches(0.3),
         font_size=14, bold=True, color=HIGHLIGHT_YELLOW)
add_multiline(slide, [
    "Social: Mất khung VIP, hiệu ứng tung XX → bạn bè thấy",
    "Gameplay: Mất DKXX boost + x3/x4 EXP → chơi tệ hẳn đi",
    "Economy: Mất Bonus %G + Upgrade R → nạp đắt hơn",
], Inches(7.1), Inches(5.6), Inches(5.4), Inches(1.0),
              font_size=11, color=TEXT_WHITE, line_spacing=1.4)

# ==============================================================
# SLIDE 12: VIP Level Benefit Table (config detail)
# ==============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "VIP Level — Config Chi Tiết")
add_text(slide, "Giá: 100k/10 ngày (tất cả level) — Benefit tích lũy cộng dồn — Chia đều 10 ngày",
         Inches(0.6), Inches(1.05), Inches(12), Inches(0.3),
         font_size=12, color=TEXT_MUTED)

headers = ["Lv", "GEM", "R.Vàng", "Chìa\nkhóa", "R.Tím", "EXP", "R.Cam",
           "Cosmetic", "DKXX", "Upgrade R", "Bonus\n%G"]
rows = [
    ["1", "1,500G", "50", "10", "5", "—", "5",
     "Khung+FX Vip1", "—", "—", "+5%"],
    ["2", "1,500G", "50", "10", "10", "x2", "5",
     "Khung+FX Vip2", "+3%", "—", "+10%"],
    ["3", "1,500G", "50", "10", "15", "x3", "10",
     "Khung+FX Vip3", "+5%", "Unlock", "+15%"],
    ["4", "1,500G", "50", "10", "20", "x3", "10",
     "Khung+FX Vip4", "—", "+10% rate", "+15%"],
    ["5", "1,500G", "50", "10", "30", "x4", "15",
     "Frame lobby", "+10%", "+15% rate", "+25%"],
]
add_table(slide, rows, headers,
          Inches(0.3), Inches(1.5), Inches(12.7),
          col_widths=[0.4, 0.9, 0.8, 0.6, 0.6, 0.5, 0.7, 1.6, 0.7, 1.0, 0.7],
          font_size=9)

# 3 benefit groups
add_shape(slide, Inches(0.5), Inches(4.0), Inches(3.8), Inches(3.0), PANEL_BG)
add_text(slide, "Resource (giữ nguyên)", Inches(0.8), Inches(4.1), Inches(3.2), Inches(0.3),
         font_size=13, bold=True, color=ACCENT2)
add_multiline(slide, [
    "1,500 GEM (750 ngay + 83G/ngày)",
    "50 rương vàng (5/ngày)",
    "10 chìa khóa đục lỗ",
    "R.Tím: 5→10→15→20→30 (tăng dần)",
    "R.Cam: 5→5→10→10→15 (tăng dần)",
], Inches(0.8), Inches(4.5), Inches(3.2), Inches(2.2),
              font_size=11, color=TEXT_WHITE, line_spacing=1.4)

add_shape(slide, Inches(4.7), Inches(4.0), Inches(3.8), Inches(3.0), PANEL_BG)
add_text(slide, "Progression (tăng theo Lv)", Inches(5.0), Inches(4.1), Inches(3.2), Inches(0.3),
         font_size=13, bold=True, color=HIGHLIGHT_YELLOW)
add_multiline(slide, [
    "EXP: — → x2 (Lv2) → x3 → x4 (Lv5)",
    "DKXX: +3% (Lv2) → +5% → +10% (Lv5)",
    "Upgrade R: Lv3 unlock → +15% (Lv5)",
    "Bonus %G: +5% → +25% (Lv5)",
    "→ Mất = mất tất cả buff",
], Inches(5.0), Inches(4.5), Inches(3.2), Inches(2.2),
              font_size=11, color=TEXT_WHITE, line_spacing=1.4)

add_shape(slide, Inches(8.9), Inches(4.0), Inches(3.8), Inches(3.0), PANEL_BG)
add_text(slide, "Social (visible)", Inches(9.2), Inches(4.1), Inches(3.2), Inches(0.3),
         font_size=13, bold=True, color=ORANGE)
add_multiline(slide, [
    "Khung avatar thay đổi theo Lv",
    "FX tung xúc xắc theo Lv",
    "Lv5: Frame lobby exclusive",
    "→ Mất VIP = bạn bè thấy",
    "→ Social pressure + loss",
], Inches(9.2), Inches(4.5), Inches(3.2), Inches(2.2),
              font_size=11, color=TEXT_WHITE, line_spacing=1.4)

# ==============================================================
# SLIDE 13: Bonus %G Migration + Upgrade R
# ==============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "Key Moves: Bonus %G + Upgrade R")

# Bonus G
add_shape(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.5), CARD_BG)
add_text(slide, "Bonus %G: Shop Trung → VIP Exclusive", Inches(0.8), Inches(1.4), Inches(5.4), Inches(0.4),
         font_size=16, bold=True, color=ACCENT)

add_table(slide, [
    ["Hiện tại", "Shop trung", "Tất cả user", "Max 20%"],
    ["Đề xuất", "VIP exclusive", "Chỉ VIP", "5%→15%"],
], ["", "Nguồn", "Ai hưởng", "Rate"],
    Inches(0.8), Inches(2.0), Inches(5.4), col_widths=[1.2, 1.3, 1.3, 1.0], font_size=10)

add_multiline(slide, [
    "Tại sao move này mạnh?",
    "",
    "1. User đang hưởng free → phải mua VIP để giữ",
    "   → Loss aversion NGAY LẬP TỨC",
    "",
    "2. Whale nạp 10,000G/tháng:",
    "   Bonus 15% = 1,500G = 1.5M VNĐ value",
    "   VIP 100k×3 = 300k → ROI cực cao",
    "",
    "3. Remove shop = giảm free GEM inflation",
    "   → Economy healthier",
], Inches(0.8), Inches(3.3), Inches(5.4), Inches(3.2),
              font_size=11, color=TEXT_DARK, line_spacing=1.2, bold_first=True)

# Upgrade R
add_shape(slide, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.5), CARD_BG)
add_text(slide, "Upgrade R: Mở thêm đường qua VIP", Inches(7.1), Inches(1.4), Inches(5.4), Inches(0.4),
         font_size=16, bold=True, color=ACCENT2)

add_multiline(slide, [
    "Hiện tại:",
    "Chỉ CLB hạng B+ mới upgrade R",
    "→ Chỉ core/old user đạt được",
    "",
    "Đề xuất:",
    "VIP Lv3+ cũng unlock Upgrade R",
    "",
    "Impact:",
    "• Mid user có thẻ R nhưng chưa CLB B → mua VIP",
    "• KHÔNG gating mới → MỞ THÊM ĐƯỜNG",
    "• Không backlash, thêm lý do cho mid segment",
    "",
    "Kết hợp Bonus %G + Upgrade R:",
    "→ VIP trở thành \"membership\" toàn diện",
    "   không chỉ subscription resource",
], Inches(7.1), Inches(2.0), Inches(5.4), Inches(4.5),
              font_size=11, color=TEXT_DARK, line_spacing=1.2, bold_first=True)

# ==============================================================
# SLIDE 14: Next Steps
# ==============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = BG_DARK

add_text(slide, "Next Steps", Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         font_size=36, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

steps = [
    "1. Chốt benefit table (review DKXX balance, Upgrade R rate)",
    "2. Estimate cost: Bonus %G cho away bao nhiêu GEM/tháng?",
    "3. Plan communication cho Bonus G migration",
    "4. Design UI/UX: countdown timer, level screen, noti flow",
    "5. A/B test trên segment nhỏ trước khi roll out",
    "6. Monitor: repeat rate, ARPT, PU, churn sau reset",
]
add_multiline(slide, steps, Inches(2), Inches(2.5), Inches(9), Inches(4),
              font_size=18, color=TEXT_WHITE, line_spacing=1.8)

add_shape(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), ACCENT)

# ==============================================================
output = "data/CTP_VIP_Full_Deck_v7.pptx"
prs.save(output)
print(f"Saved: {output}")
