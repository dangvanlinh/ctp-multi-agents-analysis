"""Generate VIP Level System Design Doc PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BG_DARK = RGBColor(0x1E, 0x27, 0x61)
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF7)
ACCENT = RGBColor(0xE6, 0x3E, 0x31)
ACCENT2 = RGBColor(0x2E, 0xA0, 0x6A)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x1E, 0x1E, 0x2E)
TEXT_MUTED = RGBColor(0x6B, 0x72, 0x80)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_HEADER = RGBColor(0x1E, 0x27, 0x61)
TABLE_ALT = RGBColor(0xEE, 0xEF, 0xF5)
HIGHLIGHT = RGBColor(0xFF, 0xD7, 0x00)
PANEL_BG = RGBColor(0x28, 0x33, 0x78)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)


def add_shape(sl, l, t, w, h, c):
    s = sl.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()

def set_cell(c, t, fs=10, b=False, col=TEXT_DARK, a=PP_ALIGN.LEFT, f=None):
    c.text = str(t); p = c.text_frame.paragraphs[0]
    p.font.size = Pt(fs); p.font.bold = b; p.font.color.rgb = col; p.alignment = a
    c.vertical_anchor = MSO_ANCHOR.MIDDLE
    if f: c.fill.solid(); c.fill.fore_color.rgb = f

def add_table(sl, rows, hdrs, l, t, w, cw=None, fs=9):
    nr, nc = len(rows)+1, len(hdrs)
    ts = sl.shapes.add_table(nr, nc, l, t, w, Inches(0.3*nr)); tb = ts.table
    if cw:
        for i, x in enumerate(cw): tb.columns[i].width = Inches(x)
    for i, h in enumerate(hdrs):
        set_cell(tb.cell(0,i), h, fs, True, TEXT_WHITE, PP_ALIGN.CENTER, TABLE_HEADER)
    for r, row in enumerate(rows):
        fl = TABLE_ALT if r%2==0 else CARD_BG
        for c, v in enumerate(row):
            al = PP_ALIGN.LEFT if c==0 else PP_ALIGN.CENTER
            set_cell(tb.cell(r+1,c), v, fs, False, TEXT_DARK, al, fl)
    return ts

def add_text(sl, t, l, tp, w, h, fs=14, b=False, c=TEXT_DARK, a=PP_ALIGN.LEFT):
    tx = sl.shapes.add_textbox(l, tp, w, h); tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = t; p.font.size = Pt(fs); p.font.bold = b
    p.font.color.rgb = c; p.font.name = "Calibri"; p.alignment = a
    return tx

def add_ml(sl, lines, l, t, w, h, fs=12, c=TEXT_DARK, ls=1.3):
    tx = sl.shapes.add_textbox(l, t, w, h); tf = tx.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = ln; p.font.size = Pt(fs); p.font.color.rgb = c; p.font.name = "Calibri"
        p.space_after = Pt(fs*(ls-1))

def title_bar(sl, t):
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = BG_LIGHT
    add_shape(sl, Inches(0), Inches(0), Inches(13.333), Inches(1.0), BG_DARK)
    add_text(sl, t, Inches(0.6), Inches(0.2), Inches(12), Inches(0.6), 24, True, TEXT_WHITE)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── SLIDE 1: Title ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = BG_DARK
add_shape(sl, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT)
add_text(sl, "VIP LEVEL SYSTEM", Inches(1), Inches(2.0), Inches(11), Inches(1),
         42, True, TEXT_WHITE, PP_ALIGN.CENTER)
add_text(sl, "Design Document", Inches(1), Inches(2.8), Inches(11), Inches(0.6),
         28, False, HIGHLIGHT, PP_ALIGN.CENTER)
add_text(sl, "Draft — 2026-04-06 — Cần review trước production", Inches(1), Inches(3.8), Inches(11), Inches(0.5),
         16, False, TEXT_WHITE, PP_ALIGN.CENTER)
add_shape(sl, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), ACCENT)

# ── SLIDE 2: Core Mechanic ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(sl, "1. Core Mechanic")
add_ml(sl, [
    "VIP có 5 cấp (Lv1-Lv5). Giá: 100k/10 ngày tất cả level.",
    "",
    "  Mua liên tiếp = +1 Level",
    "  Không mua + hết 3 ngày grace = Reset về Lv0",
    "  Chỉ thấy benefit level tiếp theo +1 (tạo curiosity)",
    "",
    "Core Pillars:",
    "  1. Fix repay trong tháng (ARPT 1→3) + repay tháng sau (17%→50%)",
    "  2. Không cắn nhau với event (SKB, CHĐ, Trứng, Lazer)",
], Inches(0.6), Inches(1.3), Inches(5.5), Inches(5.0), 14, TEXT_DARK, 1.5)

add_shape(sl, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.5), PANEL_BG)
add_text(sl, "Loss Aversion — 3 chiều", Inches(7.1), Inches(1.4), Inches(5.4), Inches(0.4),
         16, True, HIGHLIGHT)
add_ml(sl, [
    "Social: Mất khung VIP, FX tung XX → bạn bè thấy",
    "",
    "Gameplay: Mất DKXX boost + x3/x4 EXP",
    "→ Chơi tệ hẳn đi",
    "",
    "Economy: Mất Bonus %G + Upgrade R",
    "→ Nạp đắt hơn, upgrade chậm hơn",
    "",
    "→ Càng level cao càng đau",
    "→ Không ai muốn reset",
], Inches(7.1), Inches(2.0), Inches(5.4), Inches(4.0), 13, TEXT_WHITE, 1.4)

# ── SLIDE 3: Config Benefit ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(sl, "2. Config Benefit — Lv1 đến Lv5")
add_text(sl, "Giá: 100k/10 ngày tất cả level — Benefit tích lũy cộng dồn",
         Inches(0.6), Inches(1.05), Inches(12), Inches(0.3), 12, False, TEXT_MUTED)
add_table(sl, [
    ["1", "1,500G", "50", "10", "5", "—", "5", "Khung+FX Vip1", "—", "—", "TBD"],
    ["2", "1,500G", "50", "10", "10", "x2", "5", "Khung+FX Vip2", "+3%", "—", "TBD"],
    ["3", "1,500G", "50", "10", "15", "x3", "10", "Khung+FX Vip3", "+5%", "Unlock", "TBD"],
    ["4", "1,500G", "50", "10", "20", "x3", "10", "Khung+FX Vip4", "+8%", "+10%", "TBD"],
    ["5", "1,500G", "50", "10", "30", "x4", "15", "Frame lobby", "+10%", "+15%", "TBD"],
], ["Lv", "GEM", "R.Vàng", "Chìa", "R.Tím", "EXP", "R.Cam", "Cosmetic", "DKXX", "Upg R", "Bonus%G"],
    Inches(0.3), Inches(1.5), Inches(12.7),
    [0.4, 0.9, 0.7, 0.6, 0.6, 0.5, 0.6, 1.6, 0.6, 0.8, 0.7], 9)

add_text(sl, "Bonus %G: Chưa chốt — cần giải quyết migration từ shop trung",
         Inches(0.6), Inches(3.6), Inches(12), Inches(0.3), 11, True, ACCENT)

# Benefit groups
add_shape(sl, Inches(0.5), Inches(4.2), Inches(3.8), Inches(2.8), PANEL_BG)
add_text(sl, "Daily Claim (phải login)", Inches(0.8), Inches(4.3), Inches(3.2), Inches(0.3), 13, True, ACCENT2)
add_ml(sl, [
    "GEM 83G/ngày",
    "Rương vàng 5/ngày",
    "Rương tím: chia đều 10 ngày",
    "Rương cam: 5 ngày cuối (ngày 6-10)",
    "Chìa khóa: chia đều 10 ngày",
    "→ Không login = mất quà ngày đó",
], Inches(0.8), Inches(4.7), Inches(3.2), Inches(2.0), 11, TEXT_WHITE, 1.3)

add_shape(sl, Inches(4.7), Inches(4.2), Inches(3.8), Inches(2.8), PANEL_BG)
add_text(sl, "Auto-Active (có ngay)", Inches(5.0), Inches(4.3), Inches(3.2), Inches(0.3), 13, True, HIGHLIGHT)
add_ml(sl, [
    "x2/x3/x4 EXP",
    "DKXX boost",
    "Upgrade R access",
    "Cosmetic (khung, FX)",
    "Bonus %G (nếu chốt)",
    "→ Mua VIP = active, không cần claim",
], Inches(5.0), Inches(4.7), Inches(3.2), Inches(2.0), 11, TEXT_WHITE, 1.3)

add_shape(sl, Inches(8.9), Inches(4.2), Inches(3.8), Inches(2.8), PANEL_BG)
add_text(sl, "Lý do chia 2 loại", Inches(9.2), Inches(4.3), Inches(3.2), Inches(0.3), 13, True, ORANGE)
add_ml(sl, [
    "New user: daily claim → login",
    "hàng ngày → tăng retention",
    "",
    "Whale: auto-active → không phiền",
    "chỉ cần buff khi chơi",
    "",
    "2 segment, 2 behavior target",
], Inches(9.2), Inches(4.7), Inches(3.2), Inches(2.0), 11, TEXT_WHITE, 1.3)

# ── SLIDE 4: Timing & Purchase Rules ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(sl, "3. Timing, Countdown & Purchase Rules")

add_shape(sl, Inches(0.5), Inches(1.3), Inches(6.0), Inches(2.5), CARD_BG)
add_text(sl, "Timing", Inches(0.8), Inches(1.4), Inches(5.4), Inches(0.3), 16, True, BG_DARK)
add_ml(sl, [
    "VIP tính theo ngày, không theo giờ",
    "Mua ngày nào → hết 00:00 ngày thứ 11",
    "Daily claim reset lúc 00:00",
    "",
    "Grace period: 3 ngày sau hết VIP",
    "Trong grace: không quà, không buff",
    "Hết grace: reset Lv0",
], Inches(0.8), Inches(1.8), Inches(5.4), Inches(1.8), 12, TEXT_DARK, 1.3)

add_shape(sl, Inches(6.8), Inches(1.3), Inches(6.0), Inches(2.5), CARD_BG)
add_text(sl, "Mua trước hạn", Inches(7.1), Inches(1.4), Inches(5.4), Inches(0.3), 16, True, ACCENT2)
add_ml(sl, [
    "User mua tiếp khi VIP chưa hết:",
    "1. Nhận trọn quà còn lại vào inbox ngay",
    "2. Level hiện tại kết thúc",
    "3. Level mới kích hoạt ngay, 10 ngày mới",
    "",
    "→ Không mất gì, kích mua sớm = tăng ARPT",
], Inches(7.1), Inches(1.8), Inches(5.4), Inches(1.8), 12, TEXT_DARK, 1.3)

add_shape(sl, Inches(0.5), Inches(4.2), Inches(6.0), Inches(2.8), CARD_BG)
add_text(sl, "Mua nhiều lần liền", Inches(0.8), Inches(4.3), Inches(5.4), Inches(0.3), 16, True, ACCENT)
add_ml(sl, [
    "Không cho phép mua nhiều lần để nhảy level",
    "Bắt buộc 1 lần mua = 1 level",
    "",
    "Nhảy level chỉ cho:",
    "  New user offer (lần đầu → Lv2)",
    "  Lapsed offer (nếu implement sau)",
], Inches(0.8), Inches(4.7), Inches(5.4), Inches(2.0), 12, TEXT_DARK, 1.3)

add_shape(sl, Inches(6.8), Inches(4.2), Inches(6.0), Inches(2.8), CARD_BG)
add_text(sl, "Reset Rule", Inches(7.1), Inches(4.3), Inches(5.4), Inches(0.3), 16, True, ACCENT)
add_ml(sl, [
    "Hết grace 3 ngày → Reset Lv0 (mất hết)",
    "",
    "Tại sao Lv0 thay vì Lv3?",
    "  3 ngày đủ dài để quyết định",
    "  Nếu drop thì mềm hơn cũng không giữ",
    "  Loss aversion = core mechanic",
    "  Có lapsed offer để win back sau (nếu implement)",
], Inches(7.1), Inches(4.7), Inches(5.4), Inches(2.0), 12, TEXT_DARK, 1.3)

# ── SLIDE 5: Nút Kích Hoạt Logic ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(sl, "4. Nút Kích Hoạt — State Machine")

add_table(sl, [
    ["Chưa là VIP", "Không có", "Hiện (mua lần đầu)"],
    ["Đang VIP, >3 ngày, chưa claim", "Hiện", "Ẩn"],
    ["Đang VIP, >3 ngày, đã claim", "Không", "Ẩn"],
    ["Đang VIP, ≤3 ngày, chưa claim", "Hiện", "Ẩn (chờ claim)"],
    ["Đang VIP, ≤3 ngày, đã claim", "Không", "Hiện (renew)"],
    ["Hết VIP, trong 3 ngày grace", "Không", "Hiện (urgent)"],
    ["Hết VIP, quá 3 ngày → Lv0", "Không", "Hiện (mua lại Lv1)"],
], ["Trạng thái", "Nút Claim", "Nút Kích hoạt"],
    Inches(0.5), Inches(1.3), Inches(12.3),
    [5.0, 3.0, 3.5], 11)

add_shape(sl, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.5), CARD_BG)
add_text(sl, "Rule quan trọng:", Inches(0.8), Inches(4.4), Inches(11.7), Inches(0.3), 14, True, ACCENT)
add_ml(sl, [
    "Chưa claim quà hôm nay → chưa cho renew",
    "→ Đảm bảo user luôn nhận quà trước khi quyết định mua tiếp",
    "→ Daily claim là touchpoint tự nhiên để nhắc countdown VIP",
    "",
    "Không push noti — dùng in-game claim UI khi login",
], Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.5), 13, TEXT_DARK, 1.4)

# ── SLIDE 6: Upgrade R ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(sl, "5. Upgrade R — Rules")

add_shape(sl, Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5), CARD_BG)
add_ml(sl, [
    "Hiện tại: Chỉ CLB hạng B+ mới upgrade R",
    "",
    "Thêm: VIP Lv3+ cũng unlock Upgrade R",
    "→ Mở thêm đường, không gating mới",
    "→ Mid user có thẻ R nhưng chưa CLB B → mua VIP",
    "",
    "─── Edge Cases ───",
    "",
    "VIP hết → Khóa ngay quyền Upgrade R",
    "  (bấm nút instant, không có process dở dang)",
    "",
    "Thẻ R đã upgrade xong → Giữ nguyên",
    "  (không revert, chỉ mất quyền upgrade thêm)",
    "",
    "CLB B + VIP cùng lúc → Cộng dồn bonus rate",
    "  (không punish user đã đạt CLB B)",
], Inches(0.8), Inches(1.5), Inches(5.2), Inches(5.0), 13, TEXT_DARK, 1.3)

add_shape(sl, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.5), PANEL_BG)
add_text(sl, "Upgrade R theo VIP Level", Inches(7.1), Inches(1.4), Inches(5.4), Inches(0.4), 16, True, HIGHLIGHT)
add_table(sl, [
    ["Lv1-2", "—", "Không có"],
    ["Lv3", "Unlock", "Mở quyền upgrade R"],
    ["Lv4", "+10% rate", "Tăng tỷ lệ thành công"],
    ["Lv5", "+15% rate", "Tỷ lệ cao nhất"],
], ["Level", "Upgrade R", "Mô tả"],
    Inches(7.0), Inches(2.0), Inches(5.6), [0.8, 1.2, 2.8], 11)

add_text(sl, "CLB B+ có sẵn upgrade R, VIP bonus rate cộng dồn thêm",
         Inches(7.0), Inches(3.8), Inches(5.6), Inches(0.4), 11, False, TEXT_WHITE)

# ── SLIDE 7: Parked Items ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(sl, "6. Chưa Chốt — Cần Design Thêm")

# Bonus %G
add_shape(sl, Inches(0.5), Inches(1.3), Inches(4.0), Inches(5.5), CARD_BG)
add_text(sl, "Bonus %G", Inches(0.8), Inches(1.4), Inches(3.4), Inches(0.3), 16, True, ACCENT)
add_text(sl, "PARKED", Inches(2.8), Inches(1.4), Inches(1.0), Inches(0.3), 10, True, ACCENT)
add_ml(sl, [
    "Shop trung hiện: 5-20% free theo gói",
    "65% GEM payer chưa mua VIP",
    "81% transactions ≤200k (bonus ≤10%)",
    "",
    "Hướng xem xét:",
    "Giữ shop 10% flat",
    "VIP +5%/level on top",
    "",
    "Blocker: balance, GEM inflation",
], Inches(0.8), Inches(1.9), Inches(3.4), Inches(4.5), 11, TEXT_DARK, 1.3)

# Lapsed
add_shape(sl, Inches(4.8), Inches(1.3), Inches(4.0), Inches(5.5), CARD_BG)
add_text(sl, "Lapsed Rule", Inches(5.1), Inches(1.4), Inches(3.4), Inches(0.3), 16, True, ORANGE)
add_text(sl, "PARKED", Inches(7.1), Inches(1.4), Inches(1.0), Inches(0.3), 10, True, ORANGE)
add_ml(sl, [
    "Offer skip Lv2 cho lapsed user",
    "có thể bị gaming:",
    "  Bỏ VIP → chờ KM → mua khi",
    "  cần x2 EXP → loop",
    "",
    "Chờ data live để design",
    "",
    "Default: lapsed mua lại = Lv1",
    "Không có KM",
], Inches(5.1), Inches(1.9), Inches(3.4), Inches(4.5), 11, TEXT_DARK, 1.3)

# New user
add_shape(sl, Inches(9.1), Inches(1.3), Inches(3.7), Inches(5.5), CARD_BG)
add_text(sl, "New User Offer", Inches(9.4), Inches(1.4), Inches(3.1), Inches(0.3), 16, True, ACCENT2)
add_text(sl, "PARKED", Inches(11.2), Inches(1.4), Inches(1.0), Inches(0.3), 10, True, ACCENT2)
add_ml(sl, [
    "Lần mua đầu tiên trong đời",
    "→ Skip lên Lv2 ngay",
    "",
    "Cần design:",
    "  Offer UI/mockup",
    "  Pricing (100k hay KM?)",
    "",
    "Park: anh design mockup",
], Inches(9.4), Inches(1.9), Inches(3.1), Inches(4.5), 11, TEXT_DARK, 1.3)

# ── SLIDE 8: UI & Art Tasks ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(sl, "7. UI Screens & Art Tasks")

add_shape(sl, Inches(0.5), Inches(1.3), Inches(6.0), Inches(3.0), CARD_BG)
add_text(sl, "Screens đã mockup", Inches(0.8), Inches(1.4), Inches(5.4), Inches(0.3), 14, True, ACCENT2)
add_ml(sl, [
    "1. Chưa VIP — preview Lv1, nút Kích hoạt",
    "2. Đang VIP — benefit sáng, nút Claim, countdown, arrows",
    "3. VIP ≤3 ngày — urgent, Claim + Renew (disabled)",
    "4. Hết VIP — 3 ngày grace, tất cả khóa, urgent",
    "",
    "File: vip_ui_mockup.html",
], Inches(0.8), Inches(1.8), Inches(5.4), Inches(2.2), 12, TEXT_DARK, 1.3)

add_shape(sl, Inches(6.8), Inches(1.3), Inches(6.0), Inches(3.0), CARD_BG)
add_text(sl, "Arrow Navigation", Inches(7.1), Inches(1.4), Inches(5.4), Inches(0.3), 14, True, BG_DARK)
add_ml(sl, [
    "Trái: xem level đã qua",
    "Phải: level hiện tại +1 → tiếp theo \"???\"",
    "Lv5 phải: \"Max Level\"",
    "",
    "User trải nghiệm progression từ từ,",
    "không spoil level xa",
], Inches(7.1), Inches(1.8), Inches(5.4), Inches(2.2), 12, TEXT_DARK, 1.3)

add_shape(sl, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.5), CARD_BG)
add_text(sl, "Art & UI Tasks — cần design thêm", Inches(0.8), Inches(4.7), Inches(11.7), Inches(0.3), 14, True, ACCENT)
add_ml(sl, [
    "[ ] Redesign nút VIP ở bottom bar main screen (thêm level + countdown badge)",
    "[ ] New user offer popup",
    "[ ] Ref art thẻ VIP các cấp (Lv1-5), UI thay đổi theo level",
    "[ ] Ref khung quà trong UI VIP",
    "[ ] Ref cosmetic: khung avatar, FX tung XX, khung deco NV lobby, BG lobby",
], Inches(0.8), Inches(5.1), Inches(11.7), Inches(1.8), 12, TEXT_DARK, 1.3)

# ── SLIDE 9: KPI ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(sl, "8. KPI Target — Rev 350tr+/tháng")

add_text(sl, "Model: 800 PU mới, new retain 50%, old retain 80%, new ARPT=2, old ARPT=3",
         Inches(0.6), Inches(1.1), Inches(12), Inches(0.3), 13, False, TEXT_MUTED)

add_table(sl, [
    ["T1", "800", "362", "1,162", "2,686", "269M"],
    ["T2", "800", "690", "1,490", "3,670", "367M"],
    ["T3", "800", "952", "1,752", "4,456", "446M"],
    ["T4", "800", "1,162", "1,962", "5,086", "509M"],
    ["T6+", "800", "~2,000", "~2,800", "~7,600", "~760M"],
], ["Tháng", "PU mới", "Old retained", "Active PU", "Lượt", "Rev"],
    Inches(0.5), Inches(1.6), Inches(7.0), [0.8, 1.0, 1.2, 1.0, 1.0, 1.0], 11)

add_shape(sl, Inches(8.0), Inches(1.6), Inches(4.8), Inches(2.0), PANEL_BG)
add_text(sl, "Đạt 350M ngay T2", Inches(8.3), Inches(1.7), Inches(4.2), Inches(0.3), 18, True, HIGHLIGHT)
add_ml(sl, [
    "Old retain 80% → pool tích lũy",
    "Converge ~2,000 active VIP users",
    "VIP: sản phẩm phụ → revenue pillar",
], Inches(8.3), Inches(2.1), Inches(4.2), Inches(1.2), 12, TEXT_WHITE, 1.3)

add_table(sl, [
    ["New VIP PU", "800/tháng", "T3: 594 + rương cam hook"],
    ["New retain", "50%", "Hiện 17%, target gấp 3x"],
    ["Old retain", "80%", "VIP Level sunk cost giữ chân"],
    ["New ARPT", "2", "VIP 10 ngày, mua 2 = 20 ngày"],
    ["Old ARPT", "3", "Streak kích mua cả tháng"],
], ["KPI", "Target", "Cơ sở"],
    Inches(0.5), Inches(4.2), Inches(12.3), [2.0, 2.0, 6.0], 11)

# ── SLIDE 10: Summary ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = BG_DARK
add_text(sl, "Summary", Inches(1), Inches(1.0), Inches(11), Inches(0.8), 36, True, TEXT_WHITE, PP_ALIGN.CENTER)

add_ml(sl, [
    "14 decisions chốt — sẵn sàng cho dev review",
    "",
    "Đã chốt: Core mechanic, config benefit, timing, purchase rules,",
    "nút kích hoạt logic, benefit delivery, Upgrade R, UI mockup",
    "",
    "Chưa chốt: Bonus %G migration, Lapsed rule, New user offer",
    "",
    "Cần thêm: Art refs, bottom bar redesign, new user offer UI",
], Inches(2), Inches(2.2), Inches(9), Inches(4), 18, TEXT_WHITE, 1.8)

add_shape(sl, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), ACCENT)

out = "data/VIP_Design_Doc_v1.pptx"
prs.save(out)
print(f"Saved: {out}")
