"""Generate VIP Analysis PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Color palette — Cherry Bold + Dark
BG_DARK = RGBColor(0x1E, 0x27, 0x61)      # navy
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF7)     # off-white
ACCENT = RGBColor(0xE6, 0x3E, 0x31)        # cherry red
ACCENT2 = RGBColor(0x2E, 0xA0, 0x6A)       # green
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x1E, 0x1E, 0x2E)
TEXT_MUTED = RGBColor(0x6B, 0x72, 0x80)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_HEADER = RGBColor(0x1E, 0x27, 0x61)
TABLE_ALT = RGBColor(0xEE, 0xEF, 0xF5)
HIGHLIGHT_YELLOW = RGBColor(0xFF, 0xD7, 0x00)
PANEL_BG = RGBColor(0x28, 0x33, 0x78)


def add_shape(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def set_cell(cell, text, font_size=10, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT, fill=None):
    cell.text = str(text)
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill


def add_table(slide, rows_data, headers, left, top, width, col_widths=None, font_size=9):
    n_rows = len(rows_data) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.3 * n_rows))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for i, h in enumerate(headers):
        set_cell(table.cell(0, i), h, font_size=font_size, bold=True,
                 color=TEXT_WHITE, align=PP_ALIGN.CENTER, fill=TABLE_HEADER)
    for r, row in enumerate(rows_data):
        fill = TABLE_ALT if r % 2 == 0 else CARD_BG
        for c, val in enumerate(row):
            align = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            set_cell(table.cell(r + 1, c), val, font_size=font_size,
                     color=TEXT_DARK, align=align, fill=fill)
    return table_shape


def add_text(slide, text, left, top, width, height, font_size=14, bold=False,
             color=TEXT_DARK, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def title_bar(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_LIGHT
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), BG_DARK)


def bullet_card(slide, x, y, num, title, desc, col):
    add_text(slide, num, x, y, Inches(0.5), Inches(0.4), font_size=20, bold=True, color=col)
    add_text(slide, title, x + Inches(0.5), y, Inches(4.3), Inches(0.35), font_size=16, bold=True, color=TEXT_WHITE)
    add_text(slide, desc, x + Inches(0.5), y + Inches(0.35), Inches(4.3), Inches(0.3),
             font_size=11, color=RGBColor(0xAA, 0xB0, 0xCC))


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = BG_DARK
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT)
add_text(slide, "PHÂN TÍCH VIP & MONETIZATION", Inches(1), Inches(2.0), Inches(11), Inches(1),
         font_size=44, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER, font_name="Arial Black")
add_text(slide, "Cờ Tỷ Phú (CTP)", Inches(1), Inches(2.8), Inches(11), Inches(0.7),
         font_size=28, color=HIGHLIGHT_YELLOW, align=PP_ALIGN.CENTER)
add_text(slide, "Data: 01/02/2026 - 01/04/2026", Inches(1), Inches(3.8), Inches(11), Inches(0.5),
         font_size=16, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
add_shape(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), ACCENT)

# ============================================================
# SLIDE 2: Overall Revenue Breakdown
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide)
add_text(slide, "Tổng Quan Monetization — T3/2026", Inches(0.6), Inches(0.2), Inches(12), Inches(0.6),
         font_size=28, bold=True, color=TEXT_WHITE)

add_table(slide, [
    ["Tổng", "1.52 tỷ", "100%", "4,603", "330k"],
    ["GEM", "875M", "57.5%", "1,173", "746k"],
    ["Gold (gộp)", "379M", "24.9%", "3,199", "119k"],
    ["First Pay", "131M", "8.6%", "1,570", "83k"],
    ["VIP", "87M", "5.8%", "724", "121k"],
    ["SEASON_PASS", "46M", "3.0%", "387", "118k"],
], ["Sản phẩm", "Revenue", "Share", "PU", "ARPPU"],
    Inches(0.5), Inches(1.3), Inches(7.0), col_widths=[2.0, 1.2, 1.0, 1.0, 1.0], font_size=11)

add_table(slide, [
    ["Tổng", "1.16 tỷ", "1.52 tỷ", "+31%"],
    ["GEM", "888M", "875M", "-1.5%"],
    ["Gold (gộp)", "97M", "379M", "+291%"],
    ["First Pay", "72M", "131M", "+82%"],
    ["VIP", "55M", "87M", "+58%"],
    ["SEASON_PASS", "29M", "46M", "+58%"],
], ["Sản phẩm", "Rev T2", "Rev T3", "Tăng trưởng"],
    Inches(0.5), Inches(4.5), Inches(7.0), col_widths=[2.0, 1.5, 1.5, 1.5], font_size=11)

add_shape(slide, Inches(8.0), Inches(1.3), Inches(4.8), Inches(5.5), CARD_BG)
for i, (title, desc, col) in enumerate([
    ("GEM giảm 1.5%", "Dù PU tăng 16%, ARPPU giảm.\nNguồn rev chính đang yếu đi.", ACCENT),
    ("Gold/First Pay +291%", "Tăng trưởng hoàn toàn từ\nnew user do UA investment.", ACCENT2),
    ("VIP +58%", "Tăng ngang UA growth, không\noutperform. Room to improve lớn nhất.", BG_DARK),
]):
    y = Inches(1.6 + i * 1.7)
    add_text(slide, "●", Inches(8.3), y, Inches(0.3), Inches(0.3), font_size=16, bold=True, color=col)
    add_text(slide, title, Inches(8.7), y, Inches(3.8), Inches(0.35), font_size=14, bold=True, color=TEXT_DARK)
    add_text(slide, desc, Inches(8.7), y + Inches(0.4), Inches(3.8), Inches(0.7), font_size=11, color=TEXT_MUTED)

# ============================================================
# SLIDE 3: Repay Rate
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide)
add_text(slide, "Repay Rate 30 Ngày — Vấn Đề Lớn Của Game", Inches(0.6), Inches(0.2), Inches(12), Inches(0.6),
         font_size=28, bold=True, color=TEXT_WHITE)

add_shape(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(2.2), CARD_BG)
add_text(slide, "72.6%", Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.8),
         font_size=60, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, "payer mất sau 30 ngày", Inches(0.5), Inches(2.3), Inches(5.5), Inches(0.4),
         font_size=18, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
add_text(slide, "Chỉ 625 / 2,282 user T2 quay lại nạp trong 30 ngày",
         Inches(0.5), Inches(2.8), Inches(5.5), Inches(0.4), font_size=12, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

add_table(slide, [
    ["Tổng", "2,282", "625 (27.4%)", "1,657 (72.6%)"],
    ["<7d (New)", "640", "141 (22.0%)", "499 (78.0%)"],
    ["7-30d", "430", "94 (21.9%)", "336 (78.1%)"],
    ["30-90d", "376", "98 (26.1%)", "278 (73.9%)"],
    ["90-180d", "245", "75 (30.6%)", "170 (69.4%)"],
    [">180d (Core)", "591", "217 (36.7%)", "374 (63.3%)"],
], ["Segment", "Users", "Repay", "Mất"],
    Inches(7), Inches(1.3), Inches(5.8), col_widths=[1.8, 0.8, 1.5, 1.5], font_size=11)

add_shape(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.5), CARD_BG)
add_text(slide, "New user mất 78%, Core vẫn mất 63%. Cả game đều có vấn đề giữ chân payer.\n\n"
         "Game đang sống bằng UA — liên tục đổ user mới vào thay thế user cũ bỏ.\n"
         "Cần cơ chế recurring để giữ payer lâu hơn, không chỉ phụ thuộc one-time purchase.",
         Inches(1.0), Inches(4.5), Inches(11.3), Inches(1.8), font_size=14, color=TEXT_DARK)

# ============================================================
# SLIDE 4: VIP Opportunity
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide)
add_text(slide, "Tại Sao Update VIP Là Cơ Hội Lớn?", Inches(0.6), Inches(0.2), Inches(12), Inches(0.6),
         font_size=28, bold=True, color=TEXT_WHITE)

for i, (title, desc, col, num) in enumerate([
    ("VIP kém nhất\ntrong các mặt hàng",
     "VIP +58% vs Gold +291%.\nPerformance thấp nhất = room\nto improve lớn nhất.\nCải thiện cái kém dễ hơn\ncải thiện cái đang tốt.", ACCENT, "1"),
    ("Repay rate game\nchỉ 27.4%",
     "72.6% payer mất sau 30 ngày.\nVIP subscription tạo habit repay,\nkéo repay rate lên.\nLà cơ chế recurring hiếm hoi\ncủa game hiện tại.", BG_DARK, "2"),
    ("82.8% VIP\nmua 1 lần",
     "593 user vẫn chơi nhưng thành F2P.\nRedesign VIP → win back pool\nlớn nhất game.\nĐã từng chi tiền = barrier thấp\nnhất để chi lại.", ACCENT2, "3"),
]):
    x = Inches(0.5 + i * 4.2)
    add_shape(slide, x, Inches(1.3), Inches(3.8), Inches(5.5), CARD_BG)
    add_text(slide, num, x + Inches(0.15), Inches(1.5), Inches(0.6), Inches(0.6),
             font_size=28, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + Inches(0.8), Inches(1.5), Inches(2.8), Inches(0.8),
             font_size=16, bold=True, color=TEXT_DARK)
    add_text(slide, desc, x + Inches(0.3), Inches(2.6), Inches(3.2), Inches(3.5),
             font_size=12, color=TEXT_MUTED)

# ============================================================
# SLIDE 5: VIP Revenue Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide)
add_text(slide, "Tổng Quan Doanh Thu VIP", Inches(0.6), Inches(0.2), Inches(12), Inches(0.6),
         font_size=28, bold=True, color=TEXT_WHITE)

for i, (num, label, col) in enumerate([
    ("4.78%", "T2 VIP Share", TEXT_MUTED),
    ("5.75%", "T3 VIP Share", ACCENT2),
    ("+63%", "Tăng trưởng PU", ACCENT2),
    ("55.3M", "Doanh thu T2", TEXT_MUTED),
    ("87.4M", "Doanh thu T3", ACCENT),
]):
    x = Inches(0.5 + i * 2.5)
    add_shape(slide, x, Inches(1.4), Inches(2.2), Inches(1.3), CARD_BG)
    add_text(slide, num, x, Inches(1.55), Inches(2.2), Inches(0.6),
             font_size=32, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, Inches(2.1), Inches(2.2), Inches(0.3),
             font_size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

add_table(slide, [
    ["T2/2026", "1.16 tỷ", "55.3M", "4.78%", "445", "2,282"],
    ["T3/2026", "1.52 tỷ", "87.4M", "5.75%", "724", "4,603"],
], ["Tháng", "Tổng Rev", "VIP Rev", "VIP %", "VIP PU", "Tổng PU"],
    Inches(0.5), Inches(3.2), Inches(12.3), col_widths=[1.8, 2.0, 2.0, 1.5, 1.5, 1.5], font_size=11)

add_text(slide, "VIP PU tăng từ 445 (T2) lên 724 (T3), nhưng cần lưu ý: game đang đầu tư UA mạnh → A1/N1 tăng → kéo VIP PU tăng theo.\n"
         "Chưa thể kết luận VIP tự cải thiện, cần normalize theo A1/N1 để đánh giá chính xác.",
         Inches(0.5), Inches(4.8), Inches(12), Inches(0.7), font_size=13, color=TEXT_DARK)

# ============================================================
# SLIDE 6: VIP Buyers by User Age
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide)
add_text(slide, "Người Mua VIP Theo Tuổi Tài Khoản (Lần Mua Đầu)", Inches(0.6), Inches(0.2), Inches(12), Inches(0.6),
         font_size=26, bold=True, color=TEXT_WHITE)

add_table(slide, [
    ["<7 ngày (Mới)", "361", "33.1%", "42.8M", "118k", "1.17"],
    ["7-30 ngày (Sớm)", "207", "19.0%", "24.5M", "118k", "1.17"],
    ["30-90 ngày (Trung bình)", "172", "15.8%", "21.5M", "125k", "1.19"],
    ["90-180 ngày (Trung thành)", "100", "9.2%", "14.4M", "144k", "1.43"],
    [">180 ngày (Core)", "252", "23.1%", "44.4M", "176k", "1.74"],
], ["Phân khúc", "Số người", "%", "Doanh thu", "ARPPU", "TB lần mua"],
    Inches(0.5), Inches(1.4), Inches(12.3), col_widths=[3.0, 1.5, 1.0, 2.0, 1.5, 1.8], font_size=11)

for i, (num, desc, col) in enumerate([
    ("33.1%", "User mới (<7 ngày) chiếm PU nhiều nhất\nnhưng gần như không renewal (1.17 lần)", ACCENT),
    ("23.1%", "Core (>180 ngày) chỉ 23% PU nhưng\nđóng góp 30% rev, repeat cao nhất (1.74x)", ACCENT2),
    ("176k", "ARPPU tăng dần theo tuổi tài khoản:\n118k → 118k → 125k → 144k → 176k", BG_DARK),
]):
    x = Inches(0.5 + i * 4.2)
    add_shape(slide, x, Inches(4.4), Inches(3.8), Inches(2.4), CARD_BG)
    add_text(slide, num, x, Inches(4.55), Inches(3.8), Inches(0.6),
             font_size=36, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(slide, desc, x + Inches(0.3), Inches(5.2), Inches(3.2), Inches(1.2),
             font_size=11, color=TEXT_DARK, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 7: VIP Renewal/Retention
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide)
add_text(slide, "Tỷ Lệ Gia Hạn VIP — Vấn Đề Lớn Nhất", Inches(0.6), Inches(0.2), Inches(12), Inches(0.6),
         font_size=28, bold=True, color=TEXT_WHITE)

add_shape(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(2.5), CARD_BG)
add_text(slide, "82.8%", Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.9),
         font_size=60, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, "user chỉ mua VIP 1 lần rồi bỏ", Inches(0.5), Inches(2.4), Inches(5.5), Inches(0.5),
         font_size=18, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
add_text(slide, "Chỉ 17.2% mua lại lần 2+\nTrong 60 ngày, VIP 10 ngày lý tưởng nên mua 6 lần\nMax chỉ 1 user mua 10 lần",
         Inches(1.0), Inches(3.0), Inches(4.5), Inches(0.8), font_size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

add_table(slide, [
    ["1 lần", "904", "82.8%"],
    ["2 lần", "106", "9.7%"],
    ["3 lần", "34", "3.1%"],
    ["4 lần", "23", "2.1%"],
    ["5 lần", "13", "1.2%"],
    ["6+ lần", "12", "1.1%"],
], ["Số lần mua", "Số người", "%"],
    Inches(7), Inches(1.3), Inches(5.5), col_widths=[2.0, 1.5, 1.5], font_size=11)

add_shape(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.5), CARD_BG)
add_text(slide, "VIP không tạo được stickiness. Value proposition sau lần mua đầu không đủ mạnh.\n"
         "Renewal rate ~17% là cực thấp cho subscription model. Cần loyalty/streak reward để incentivize.",
         Inches(1.0), Inches(5.2), Inches(11.3), Inches(1.0), font_size=13, color=TEXT_DARK)

# ============================================================
# SLIDE 8: Cross-Spending Behavior
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide)
add_text(slide, "Hành Vi Chi Tiêu Của Người Mua VIP", Inches(0.6), Inches(0.2), Inches(12), Inches(0.6),
         font_size=28, bold=True, color=TEXT_WHITE)

for i, (num, label, col) in enumerate([
    ("62%", "Người mua VIP cũng nạp GEM", ACCENT),
    ("1.77M", "GEM ARPPU của người mua VIP", BG_DARK),
    ("9%", "VIP chỉ chiếm 9% chi tiêu\ncủa chính người mua VIP", ACCENT2),
]):
    x = Inches(0.5 + i * 4.2)
    add_shape(slide, x, Inches(1.3), Inches(3.8), Inches(1.6), CARD_BG)
    add_text(slide, num, x, Inches(1.4), Inches(3.8), Inches(0.7),
             font_size=36, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, Inches(2.1), Inches(3.8), Inches(0.6),
             font_size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

add_table(slide, [
    ["GEM", "1.205 tỷ", "680", "1.77M"],
    ["VIP", "147.6M", "1,092", "135k"],
    ["GOLD", "79.3M", "266", "298k"],
    ["BIG_RICH", "65.3M", "379", "172k"],
    ["GOLD_BONUS", "55.5M", "300", "185k"],
    ["SEASON_PASS", "44.5M", "312", "143k"],
], ["Nguồn", "Doanh thu", "Số người", "ARPPU"],
    Inches(0.5), Inches(3.3), Inches(12.3), col_widths=[3.0, 3.0, 2.0, 3.0], font_size=11)

add_text(slide, "Người mua VIP là mid-to-heavy spenders, không phải casual. VIP là sản phẩm phụ, "
         "GEM là main purchase. Cơ hội: upsell VIP cho nhóm đã sẵn sàng chi tiêu.",
         Inches(0.5), Inches(6.0), Inches(12), Inches(0.5), font_size=13, color=TEXT_DARK)

# ============================================================
# SLIDE 9: GEM Spending by VIP Buyer Age
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide)
add_text(slide, "Phân Bổ GEM Spending Theo Tuổi Tài Khoản (680 VIP Buyers)", Inches(0.6), Inches(0.2), Inches(12), Inches(0.6),
         font_size=26, bold=True, color=TEXT_WHITE)

add_table(slide, [
    ["<7 ngày (Mới)", "211", "31%", "250M", "21%", "1.19M"],
    ["7-30 ngày (Sớm)", "126", "19%", "136M", "11%", "1.08M"],
    ["30-90 ngày (TB)", "114", "17%", "148M", "12%", "1.30M"],
    ["90-180 ngày (Trung thành)", "69", "10%", "77M", "6%", "1.12M"],
    [">180 ngày (Core)", "160", "24%", "593M", "49%", "3.71M"],
], ["Phân khúc", "GEM Payers", "% PU", "GEM Revenue", "% Rev", "ARPPU"],
    Inches(0.5), Inches(1.3), Inches(7.5), col_widths=[2.5, 1.1, 0.8, 1.3, 0.8, 1.0], font_size=11)

for i, (num, title, desc, col) in enumerate([
    ("49%", "GEM rev từ Core >180 ngày", "Chỉ 160 người nhưng chi 593M\nARPPU 3.71M — gấp 3x các nhóm khác\nĐây là whale thực sự của game", ACCENT),
    ("31%", "PU từ User mới <7 ngày", "211 người mua GEM, đông nhất\nnhưng ARPPU chỉ 1.19M\nMua thử nhiều, chưa phải big spender", ACCENT2),
    ("1.1-1.3M", "ARPPU vùng giữa 7-180d", "Không segment nào nổi bật\nConfirm churn cliff sau tháng đầu\nUser không scale spending theo thời gian", BG_DARK),
]):
    y = Inches(1.3 + i * 1.9)
    add_shape(slide, Inches(8.5), y, Inches(4.3), Inches(1.7), CARD_BG)
    add_text(slide, num, Inches(8.7), y + Inches(0.1), Inches(1.5), Inches(0.5),
             font_size=28, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(slide, title, Inches(10.2), y + Inches(0.1), Inches(2.4), Inches(0.35),
             font_size=12, bold=True, color=TEXT_DARK)
    add_text(slide, desc, Inches(10.2), y + Inches(0.5), Inches(2.4), Inches(1.0),
             font_size=10, color=TEXT_MUTED)

add_shape(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.2), CARD_BG)
add_text(slide, "Kết hợp VIP + GEM data: Core user vừa mua VIP nhiều lần (repeat 1.74x) vừa nạp GEM cực mạnh (593M). "
         "Nếu mất 160 core users này = mất ~593M GEM + phần lớn VIP recurring. Đây là nhóm cần protect nhất.",
         Inches(1.0), Inches(6.0), Inches(11.3), Inches(0.8), font_size=13, color=TEXT_DARK)

# ============================================================
# SLIDE 10: Deep-dive Funnel
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide)
add_text(slide, "Deep-dive: 82.8% Chỉ Mua VIP 1 Lần — Chuyện Gì Xảy Ra?", Inches(0.6), Inches(0.2), Inches(12), Inches(0.6),
         font_size=26, bold=True, color=TEXT_WHITE)

add_shape(slide, Inches(0.5), Inches(1.3), Inches(7.0), Inches(5.5), CARD_BG)
add_text(slide, "Phễu sau khi VIP hết hạn", Inches(0.8), Inches(1.4), Inches(6.4), Inches(0.4),
         font_size=16, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

for left, top, width, label, pct, color in [
    (Inches(0.8), Inches(2.1), Inches(6.4), "1,092 VIP Buyers", "100%", BG_DARK),
    (Inches(1.2), Inches(2.85), Inches(5.6), "904 mua 1 lần", "82.8%", ACCENT),
    (Inches(1.6), Inches(3.6), Inches(4.8), "~700 vẫn chơi game", "77%", RGBColor(0xE6, 0x7E, 0x22)),
    (Inches(2.0), Inches(4.35), Inches(4.0), "~593 không nạp gì", "85%", RGBColor(0xC0, 0x39, 0x2B)),
]:
    add_shape(slide, left, top, width, Inches(0.6), color)
    add_text(slide, f"{label}  ({pct})", left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.4),
             font_size=13, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "→ 593 user đang chơi, từng chi tiền,\n   giờ thành F2P hoàn toàn",
         Inches(1.5), Inches(5.2), Inches(5.5), Inches(0.6),
         font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, "Đo trong 7 ngày sau VIP hết hạn. Chỉ lấy user có VIP hết trước 25/03.",
         Inches(1.0), Inches(5.85), Inches(6.0), Inches(0.3), font_size=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

add_shape(slide, Inches(8.0), Inches(1.3), Inches(4.8), Inches(5.5), CARD_BG)
add_text(slide, "Phát hiện chính", Inches(8.3), Inches(1.4), Inches(4.2), Inches(0.4),
         font_size=16, bold=True, color=BG_DARK)

for i, (title, desc, col) in enumerate([
    ("Không phải churn game", "77% user 1-lần vẫn active sau VIP hết.\nHọ không rời game — họ rời VIP.", ACCENT2),
    ("VIP = cửa cuối trước F2P", "85% không nạp bất kỳ gì sau VIP hết.\nVIP không bridge sang spending khác.", ACCENT),
    ("Window repay = 15 ngày", "77% repeat buyers mua lại trong 15 ngày.\nBinary: subscribe ngay hoặc mất luôn.\nKhông có cơ chế kích repay hiện tại.", BG_DARK),
]):
    y = Inches(2.0 + i * 1.65)
    add_text(slide, "●", Inches(8.3), y, Inches(0.3), Inches(0.3), font_size=16, bold=True, color=col)
    add_text(slide, title, Inches(8.7), y, Inches(3.8), Inches(0.3), font_size=13, bold=True, color=TEXT_DARK)
    add_text(slide, desc, Inches(8.7), y + Inches(0.35), Inches(3.8), Inches(0.9), font_size=10, color=TEXT_MUTED)

# ============================================================
# SLIDE 11: Key Insights & Next Steps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = BG_DARK
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT)

add_text(slide, "Nhận Định Chính & Câu Hỏi Mở", Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         font_size=32, bold=True, color=TEXT_WHITE)

add_shape(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.5), PANEL_BG)
for i, (num, title, desc) in enumerate([
    ("1.", "VIP PU tăng theo UA", "+63% PU, nhưng do đầu tư UA → A1/N1 tăng"),
    ("2.", "Core >180 ngày là trụ cột", "23% PU, 30% rev, repeat 1.74x"),
    ("3.", "82.8% chỉ mua 1 lần", "77% vẫn chơi, 85% ngừng nạp → VIP là cửa cuối trước F2P"),
    ("4.", "Window kích repay = 15 ngày", "77% repeat mua lại trong 15 ngày, sau đó mất luôn"),
    ("5.", "Game repay rate chỉ 27.4%", "72.6% payer mất sau 30 ngày, VIP subscription là cơ hội"),
]):
    bullet_card(slide, Inches(0.8), Inches(1.6 + i * 1.0), num, title, desc, ACCENT)

add_shape(slide, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.5), PANEL_BG)
add_text(slide, "Câu hỏi mở", Inches(7.3), Inches(1.5), Inches(5.2), Inches(0.5),
         font_size=18, bold=True, color=HIGHLIGHT_YELLOW)

for i, q in enumerate([
    "Segment VIP thế nào để phục vụ\ncả new user lẫn core user?",
    "User mới cần benefit gì khác\nngoài x2 EXP để thấy giá trị sớm hơn?",
    "Cần loyalty/streak reward để tăng\ntỷ lệ gia hạn từ 17% lên bao nhiêu?",
    "VIP share hiện 5.75%\n→ Target bao nhiêu % là hợp lý?",
]):
    add_text(slide, q, Inches(7.5), Inches(2.2 + i * 1.15), Inches(5.0), Inches(0.9),
             font_size=12, color=TEXT_WHITE)

add_shape(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), ACCENT)

# Save
output_path = "data/VIP_Analysis_CTP_v7.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
