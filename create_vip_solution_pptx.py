"""Generate VIP Solution Concept PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Color palette
BG_DARK = RGBColor(0x1E, 0x27, 0x61)
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF7)
ACCENT = RGBColor(0xE6, 0x3E, 0x31)
ACCENT2 = RGBColor(0x2E, 0xA0, 0x6A)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x1E, 0x1E, 0x2E)
TEXT_MUTED = RGBColor(0x6B, 0x72, 0x80)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_HEADER = RGBColor(0x1E, 0x27, 0x61)
TABLE_ALT = RGBColor(0xEE, 0xEF, 0xF5)
HIGHLIGHT_YELLOW = RGBColor(0xFF, 0xD7, 0x00)
PANEL_BG = RGBColor(0x28, 0x33, 0x78)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)


def add_shape(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def set_cell(cell, text, font_size=10, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT, fill=None):
    cell.text = str(text)
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill


def add_table(slide, rows_data, headers, left, top, width, col_widths=None, font_size=9):
    n_rows = len(rows_data) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.3 * n_rows))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for i, h in enumerate(headers):
        set_cell(table.cell(0, i), h, font_size=font_size, bold=True,
                 color=TEXT_WHITE, align=PP_ALIGN.CENTER, fill=TABLE_HEADER)
    for r, row in enumerate(rows_data):
        fill = TABLE_ALT if r % 2 == 0 else CARD_BG
        for c, val in enumerate(row):
            align = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            set_cell(table.cell(r + 1, c), val, font_size=font_size,
                     color=TEXT_DARK, align=align, fill=fill)
    return table_shape


def add_text(slide, text, left, top, width, height, font_size=14, bold=False,
             color=TEXT_DARK, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_multiline(slide, lines, left, top, width, height, font_size=12,
                  color=TEXT_DARK, line_spacing=1.3, bold_first=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(font_size * (line_spacing - 1))
        if bold_first and i == 0:
            p.font.bold = True
    return txBox


def title_bar(slide, title_text):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_LIGHT
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), BG_DARK)
    add_text(slide, title_text, Inches(0.6), Inches(0.2), Inches(12), Inches(0.6),
             font_size=26, bold=True, color=TEXT_WHITE)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = BG_DARK
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT)
add_text(slide, "VIP SOLUTION CONCEPT", Inches(1), Inches(2.0), Inches(11), Inches(1),
         font_size=44, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER, font_name="Arial Black")
add_text(slide, "VIP Level System — Cờ Tỷ Phú", Inches(1), Inches(2.8), Inches(11), Inches(0.7),
         font_size=28, color=HIGHLIGHT_YELLOW, align=PP_ALIGN.CENTER)
add_text(slide, "Target: 300tr/tháng | PU 1,000 | ARPT 3 | ARPPU 300k", Inches(1), Inches(3.8), Inches(11), Inches(0.5),
         font_size=16, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
add_shape(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), ACCENT)

# ============================================================
# SLIDE 2: VIP Current Problems
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "Vấn đề hiện tại của VIP")

add_text(slide, "VIP hiện tại: 100k/10 ngày — 82.8% chỉ mua 1 lần", Inches(0.6), Inches(1.1), Inches(12), Inches(0.4),
         font_size=16, bold=True, color=ACCENT)

add_table(slide, [
    ["Rev/tháng", "~50-87M", "5.75% total rev"],
    ["PU/tháng", "487-724", "ARPPU 100k"],
    ["Repeat rate", "17.2%", "82.8% mua 1 lần rồi bỏ"],
    ["ARPT", "1.0", "Lý tưởng: 3 (mua liên tục cả tháng)"],
    ["Repay next month", "~21%", "79% không quay lại tháng sau"],
], ["Metric", "Hiện tại", "Ghi chú"],
    Inches(0.6), Inches(1.7), Inches(7), col_widths=[2.0, 1.5, 3.5], font_size=11)

# Root causes panel
add_shape(slide, Inches(8.0), Inches(1.7), Inches(4.8), Inches(5.0), PANEL_BG)
add_text(slide, "Root Causes", Inches(8.3), Inches(1.8), Inches(4.2), Inches(0.4),
         font_size=18, bold=True, color=HIGHLIGHT_YELLOW)

causes = [
    "1. Không repay trong tháng:",
    "   • VIP vs Free như nhau về hình ảnh",
    "   • x2 EXP chỉ cần khi upgrade, xong = hết nhu cầu",
    "   • Mua lần 2 không khác gì lần 1",
    "   • Không có đặc quyền exclusive",
    "",
    "2. Không repay tháng sau:",
    "   • 2 lần mua độc lập, không có streak/loyalty",
    "   • Không noti khi gần hết VIP",
    "   • Không có content mới mỗi tháng",
    "",
    "3. Không convert new pay:",
    "   • 58% payer chưa từng mua VIP",
    "   • 42% lapsed — từng mua nhưng bỏ",
]
add_multiline(slide, causes, Inches(8.3), Inches(2.3), Inches(4.2), Inches(4.2),
              font_size=11, color=TEXT_WHITE)

# ============================================================
# SLIDE 3: KPI Framework
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "KPI Framework — Từ ~50tr lên 300tr/tháng")

# KPI breakdown table
add_text(slide, "Phân rã KPI: Rev = PU × ARPT × ARPPT", Inches(0.6), Inches(1.1), Inches(12), Inches(0.4),
         font_size=16, bold=True, color=TEXT_DARK)

add_table(slide, [
    ["Rev/tháng", "~50tr", "300tr", "6x"],
    ["PU/tháng", "487", "1,000", "+513 PU"],
    ["ARPT (lần mua/user)", "1.0", "3.0", "Mua liên tục cả tháng"],
    ["ARPPU", "100k", "300k", "= ARPT × 100k"],
], ["Metric", "Hiện tại", "KPI Target", "Gap"],
    Inches(0.6), Inches(1.7), Inches(7.5), col_widths=[2.5, 1.5, 1.5, 2.0], font_size=11)

# 3 mechanics
add_text(slide, "3 Mechanics cần fix", Inches(0.6), Inches(3.6), Inches(12), Inches(0.4),
         font_size=18, bold=True, color=BG_DARK)

# Mechanic 1
add_shape(slide, Inches(0.6), Inches(4.2), Inches(3.8), Inches(2.8), PANEL_BG)
add_text(slide, "M1: Convert New Pay", Inches(0.9), Inches(4.3), Inches(3.2), Inches(0.4),
         font_size=14, bold=True, color=HIGHLIGHT_YELLOW)
add_text(slide, "KPI: 600 PU (+200)", Inches(0.9), Inches(4.7), Inches(3.2), Inches(0.3),
         font_size=12, bold=True, color=ACCENT)
m1_lines = [
    "• 58% (223) chưa mua VIP",
    "  → Offer lần đầu",
    "• 42% (159) lapsed",
    "  → Offer welcome back VIP",
]
add_multiline(slide, m1_lines, Inches(0.9), Inches(5.2), Inches(3.2), Inches(1.5),
              font_size=11, color=TEXT_WHITE)

# Mechanic 2
add_shape(slide, Inches(4.8), Inches(4.2), Inches(3.8), Inches(2.8), PANEL_BG)
add_text(slide, "M2: Repay Next Month", Inches(5.1), Inches(4.3), Inches(3.2), Inches(0.4),
         font_size=14, bold=True, color=HIGHLIGHT_YELLOW)
add_text(slide, "KPI: Repay 400 (+300 PU)", Inches(5.1), Inches(4.7), Inches(3.2), Inches(0.3),
         font_size=12, bold=True, color=ACCENT)
m2_lines = [
    "• Streak tháng (loyalty tích lũy)",
    "• Noti những ngày cuối VIP",
    "• Renew content mỗi tháng",
]
add_multiline(slide, m2_lines, Inches(5.1), Inches(5.2), Inches(3.2), Inches(1.5),
              font_size=11, color=TEXT_WHITE)

# Mechanic 3
add_shape(slide, Inches(9.0), Inches(4.2), Inches(3.8), Inches(2.8), PANEL_BG)
add_text(slide, "M3: Repay Trong Tháng", Inches(9.3), Inches(4.3), Inches(3.2), Inches(0.4),
         font_size=14, bold=True, color=HIGHLIGHT_YELLOW)
add_text(slide, "KPI: ARPT 3 (mua 3 lần/tháng)", Inches(9.3), Inches(4.7), Inches(3.2), Inches(0.3),
         font_size=12, bold=True, color=ACCENT)
m3_lines = [
    "• Cosmetic (cảm giác mất mát)",
    "• Benefit tăng dần theo level",
    "• Streak benefit (progression)",
    "• Limited access / exclusive",
]
add_multiline(slide, m3_lines, Inches(9.3), Inches(5.2), Inches(3.2), Inches(1.5),
              font_size=11, color=TEXT_WHITE)

# ============================================================
# SLIDE 4: Solution Overview — VIP Level System
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "Solution: VIP Level System")

add_text(slide, "Core Mechanic: Mỗi lần mua VIP liên tiếp = +1 Level. Max Lv5. Không mua tiếp = Reset về Lv0.",
         Inches(0.6), Inches(1.1), Inches(12), Inches(0.4),
         font_size=16, bold=True, color=ACCENT)

# Rules
add_shape(slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(5.2), CARD_BG)
add_text(slide, "Rule chung", Inches(0.9), Inches(1.8), Inches(5.2), Inches(0.4),
         font_size=16, bold=True, color=BG_DARK)
rules = [
    "• VIP có 6 cấp (Lv0-Lv5). Mỗi cấp unlock quyền lợi mới",
    "• Mỗi VIP trị giá 100k, kéo dài 10 ngày",
    "• Hết 10 ngày → countdown 1 ngày để nạp level tiếp",
    "• Nếu không nạp trong countdown → reset về Lv0",
    "• Level 5: countdown 3 ngày (ưu đãi core user)",
    "• Timeline 10 ngày đi riêng theo user",
    "• Quyền lợi ẩn: chỉ thấy level tiếp theo +1",
    "  (VD: đang Lv1 → thấy quà Lv2, Lv3+ bị ẩn)",
]
add_multiline(slide, rules, Inches(0.9), Inches(2.3), Inches(5.2), Inches(3.5),
              font_size=12, color=TEXT_DARK, line_spacing=1.5)

# Loss aversion panel
add_shape(slide, Inches(6.8), Inches(1.7), Inches(6.0), Inches(2.4), PANEL_BG)
add_text(slide, "Tại sao tạo \"loss aversion\"?", Inches(7.1), Inches(1.8), Inches(5.4), Inches(0.4),
         font_size=16, bold=True, color=HIGHLIGHT_YELLOW)
loss_lines = [
    "• Mất cosmetic (khung VIP, hiệu ứng) → bạn bè thấy",
    "• Mất DKXX boost → chơi tệ hẳn đi",
    "• Từ x4 EXP về x1 → level chậm cực kỳ",
    "• Mất Bonus %G khi nạp → nạp GEM đắt hơn",
    "• Mất quyền Upgrade R → chặn progression",
    "→ Đau ở CẢ 3 chiều: social + gameplay + economy",
]
add_multiline(slide, loss_lines, Inches(7.1), Inches(2.3), Inches(5.4), Inches(1.6),
              font_size=12, color=TEXT_WHITE, line_spacing=1.4)

# Lapsed rule
add_shape(slide, Inches(6.8), Inches(4.3), Inches(6.0), Inches(1.3), RGBColor(0x3A, 0x45, 0x8A))
add_text(slide, "Rule cho Lapsed User", Inches(7.1), Inches(4.4), Inches(5.4), Inches(0.4),
         font_size=14, bold=True, color=ORANGE)
add_multiline(slide, [
    "• Qua 1 mùa không nạp VIP → popup KM nạp lên Lv2 luôn",
    "• Áp dụng 1 lần/user (anti-gaming)",
], Inches(7.1), Inches(4.9), Inches(5.4), Inches(0.6),
              font_size=12, color=TEXT_WHITE)

# Convert new
add_shape(slide, Inches(6.8), Inches(5.8), Inches(6.0), Inches(1.1), RGBColor(0x3A, 0x45, 0x8A))
add_text(slide, "Rule cho New User", Inches(7.1), Inches(5.9), Inches(5.4), Inches(0.4),
         font_size=14, bold=True, color=ACCENT2)
add_multiline(slide, [
    "• 58% payer chưa mua VIP → Offer lần đầu (giá KM hoặc bonus)",
    "• Rương cam ở Lv1 = hook mạnh cho new user",
], Inches(7.1), Inches(6.3), Inches(5.4), Inches(0.6),
              font_size=12, color=TEXT_WHITE)

# ============================================================
# SLIDE 5: VIP Level Benefit Table (updated concept)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "VIP Level — Bảng Benefit Chi Tiết")

add_text(slide, "Giá: 100k/10 ngày (tất cả level) — Benefit tích lũy cộng dồn", Inches(0.6), Inches(1.1), Inches(12), Inches(0.4),
         font_size=14, bold=True, color=TEXT_MUTED)

headers = ["Level", "GEM", "Rương vàng", "Chìa khóa", "Rương tím", "EXP", "Rương cam",
           "Cosmetic", "DKXX", "Upgrade R", "Bonus %G"]
rows = [
    ["Lv1", "1,500G", "50", "10", "5", "x2", "5",
     "Khung + FX Vip1", "—", "—", "+5%"],
    ["Lv2", "1,500G", "50", "10", "5", "x2", "5 (random)",
     "Khung + FX Vip2", "+3%", "—", "+10%"],
    ["Lv3", "1,500G", "50", "10", "5", "x3", "5 (random)",
     "Khung + FX Vip3", "+5%", "Unlock", "+15%"],
    ["Lv4", "1,500G", "50", "10", "5", "x3", "5 (random)",
     "Khung + FX Vip4", "+5%", "+10% rate", "+15%"],
    ["Lv5", "1,500G", "50", "10", "5", "x4", "5 (random)",
     "Frame lobby Vip5", "+10%", "+15% rate", "+15%"],
]
add_table(slide, rows, headers,
          Inches(0.3), Inches(1.7), Inches(12.7),
          col_widths=[0.6, 0.9, 0.9, 0.8, 0.8, 0.5, 0.9, 1.6, 0.7, 0.9, 0.8],
          font_size=9)

# Highlight boxes
add_text(slide, "Chia đều 10 ngày", Inches(0.9), Inches(3.9), Inches(3.5), Inches(0.3),
         font_size=10, color=TEXT_MUTED)

# Benefit grouping
add_shape(slide, Inches(0.6), Inches(4.5), Inches(3.8), Inches(2.5), PANEL_BG)
add_text(slide, "Resource (giữ nguyên mọi level)", Inches(0.9), Inches(4.6), Inches(3.2), Inches(0.3),
         font_size=13, bold=True, color=ACCENT2)
add_multiline(slide, [
    "• 1,500 GEM (750 ngay + 83G/ngày × 9)",
    "• 50 rương vàng (5/ngày)",
    "• 10 chìa khóa đục lỗ trang sức",
    "• 5 rương tím",
    "• 5 rương cam (hook new user)",
], Inches(0.9), Inches(5.0), Inches(3.2), Inches(1.8),
              font_size=11, color=TEXT_WHITE)

add_shape(slide, Inches(4.8), Inches(4.5), Inches(3.8), Inches(2.5), PANEL_BG)
add_text(slide, "Progression (tăng theo level)", Inches(5.1), Inches(4.6), Inches(3.2), Inches(0.3),
         font_size=13, bold=True, color=HIGHLIGHT_YELLOW)
add_multiline(slide, [
    "• EXP: x2 → x3 (Lv3) → x4 (Lv5)",
    "• DKXX: +3% (Lv2) → +10% (Lv5)",
    "• Upgrade R: unlock (Lv3) → +15% (Lv5)",
    "• Bonus %G nạp: 5% → 15%",
    "→ Mất level = mất tất cả buff này",
], Inches(5.1), Inches(5.0), Inches(3.2), Inches(1.8),
              font_size=11, color=TEXT_WHITE)

add_shape(slide, Inches(9.0), Inches(4.5), Inches(3.8), Inches(2.5), PANEL_BG)
add_text(slide, "Social (visible cho người khác)", Inches(9.3), Inches(4.6), Inches(3.2), Inches(0.3),
         font_size=13, bold=True, color=ORANGE)
add_multiline(slide, [
    "• Khung avatar ingame thay đổi theo Lv",
    "• Hiệu ứng tung xúc xắc theo Lv",
    "• Lv5: Frame avatar lobby exclusive",
    "→ Mất VIP = bạn bè thấy mất khung",
    "→ Social pressure + loss aversion",
], Inches(9.3), Inches(5.0), Inches(3.2), Inches(1.8),
              font_size=11, color=TEXT_WHITE)

# ============================================================
# SLIDE 6: Mechanic giải quyết từng root cause
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "Benefit → Root Cause Mapping")

add_text(slide, "Mỗi benefit được thiết kế để giải quyết 1 root cause cụ thể",
         Inches(0.6), Inches(1.1), Inches(12), Inches(0.4),
         font_size=14, bold=True, color=TEXT_MUTED)

headers = ["Root Cause", "Benefit giải quyết", "Cách tác động", "KPI tăng"]
rows = [
    ["VIP vs Free như nhau hình ảnh",
     "Cosmetic (khung, FX, frame lobby)",
     "Mất VIP = mất identity trước bạn bè",
     "Repay trong tháng"],
    ["x2 EXP chỉ cần khi upgrade",
     "x2→x4 EXP + DKXX boost",
     "Buff permanent khi có VIP, mất = chơi tệ đi",
     "Repay trong tháng"],
    ["Mua lần 2 = lần 1, không progression",
     "VIP Level tăng dần + Streak benefit",
     "Mỗi lần mua unlock thêm, dừng = reset Lv0",
     "Repay trong tháng"],
    ["Không có đặc quyền exclusive",
     "Upgrade R + Bonus %G exclusive VIP",
     "Remove bonus G khỏi shop → chỉ VIP có",
     "Repay trong tháng"],
    ["Không có streak/loyalty tháng",
     "Level tích lũy cross-tháng",
     "Streak 50 ngày Lv5 = sunk cost cực lớn",
     "Repay next month"],
    ["Không noti khi gần hết VIP",
     "Countdown + push noti",
     "1 ngày countdown → urgency mua tiếp",
     "Repay next month"],
    ["Content tháng sau = tháng trước",
     "Renew content (rương cam mới, event link)",
     "Rương cam content thay đổi theo event/mùa",
     "Repay next month"],
    ["58% payer chưa mua VIP",
     "Offer lần đầu + rương cam Lv1 hook",
     "Rương cam = wow factor cho new user",
     "Convert new pay"],
    ["42% lapsed user",
     "Popup KM skip lên Lv2",
     "Giảm barrier quay lại, thấy value ngay",
     "Convert new pay"],
]
add_table(slide, rows, headers,
          Inches(0.3), Inches(1.7), Inches(12.7),
          col_widths=[3.0, 3.2, 3.5, 2.0],
          font_size=10)

# ============================================================
# SLIDE 7: Bonus %G — Migration from Shop
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "Bonus %G — Chuyển từ Shop Trung sang VIP Exclusive")

add_text(slide, "Hiện tại: Shop trung bonus max 20% → Đề xuất: Remove khỏi shop, chuyển thành VIP benefit",
         Inches(0.6), Inches(1.1), Inches(12), Inches(0.4),
         font_size=14, bold=True, color=ACCENT)

add_table(slide, [
    ["Hiện tại (Shop trung)", "Tất cả user", "Max 20%", "Miễn phí"],
    ["Đề xuất (VIP)", "Chỉ VIP user", "5% → 15% theo level", "100k/10 ngày"],
], ["", "Ai được hưởng", "Bonus rate", "Chi phí"],
    Inches(0.6), Inches(1.8), Inches(7), col_widths=[2.2, 1.8, 2.0, 1.5], font_size=11)

# Impact
add_shape(slide, Inches(0.6), Inches(3.2), Inches(5.8), Inches(3.5), PANEL_BG)
add_text(slide, "Tại sao move này mạnh?", Inches(0.9), Inches(3.3), Inches(5.2), Inches(0.4),
         font_size=16, bold=True, color=HIGHLIGHT_YELLOW)
add_multiline(slide, [
    "1. User đang hưởng bonus free → phải mua VIP để giữ lại",
    "   → Tạo loss aversion NGAY LẬP TỨC cho existing payer",
    "",
    "2. Whale nạp 10,000G/tháng, bonus 15% = 1,500G extra",
    "   → VIP 100k × 3 = 300k để được 1.5M VNĐ value",
    "   → ROI cực cao, không lý do gì không mua",
    "",
    "3. Remove từ shop = giảm free GEM inflation",
    "   → Economy healthier",
]
, Inches(0.9), Inches(3.8), Inches(5.2), Inches(2.8),
              font_size=12, color=TEXT_WHITE, line_spacing=1.3)

# Upgrade R panel
add_shape(slide, Inches(6.8), Inches(3.2), Inches(6.0), Inches(3.5), PANEL_BG)
add_text(slide, "Upgrade R — Mở thêm đường, không gating", Inches(7.1), Inches(3.3), Inches(5.4), Inches(0.4),
         font_size=16, bold=True, color=HIGHLIGHT_YELLOW)
add_multiline(slide, [
    "Hiện tại: Chỉ CLB hạng B+ mới upgrade R",
    "→ Chỉ core/old user đạt được",
    "",
    "Đề xuất: VIP Lv3+ cũng unlock Upgrade R",
    "→ Mid user có thẻ R từ event nhưng chưa CLB B",
    "   cũng upgrade được nếu mua VIP",
    "",
    "Không phải gating mới → MỞ THÊM ĐƯỜNG",
    "→ Không gây backlash, thêm lý do mua VIP cho mid",
], Inches(7.1), Inches(3.8), Inches(5.4), Inches(2.8),
              font_size=12, color=TEXT_WHITE, line_spacing=1.3)

# ============================================================
# SLIDE 8: KPI Impact Estimation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(slide, "KPI Impact Estimation")

add_text(slide, "Rev target 300tr/tháng = PU 1,000 × ARPT 3 × ARPPU 100k",
         Inches(0.6), Inches(1.1), Inches(12), Inches(0.4),
         font_size=16, bold=True, color=TEXT_DARK)

# Compound model
add_text(slide, "Compound Model (giả định 700 PU mới/tháng, repeat 50%, freq 3):",
         Inches(0.6), Inches(1.7), Inches(12), Inches(0.3),
         font_size=13, bold=True, color=BG_DARK)

add_table(slide, [
    ["T1", "700-900", "0", "700-900", "70-90M"],
    ["T2", "700-900", "350-450", "1,750-1,960", "175-196M"],
    ["T3", "700-900", "525-588", "2,275-2,464", "228-246M"],
    ["T4", "700-900", "613-655", "2,538-2,666", "254-267M"],
    ["T6+", "700-900", "700-900", "2,800-3,600", "280-360M"],
], ["Tháng", "PU mới", "Retained", "Total lượt mua", "Rev estimate"],
    Inches(0.6), Inches(2.2), Inches(7), col_widths=[1.0, 1.2, 1.2, 1.6, 1.5], font_size=11)

# Impact as % total rev
add_shape(slide, Inches(8.0), Inches(2.2), Inches(4.8), Inches(2.0), PANEL_BG)
add_text(slide, "Impact vs Total Rev", Inches(8.3), Inches(2.3), Inches(4.2), Inches(0.4),
         font_size=14, bold=True, color=HIGHLIGHT_YELLOW)
add_multiline(slide, [
    "Hiện tại: ~50-87M (5.75% total rev)",
    "Target T6: 280-360M (~20% total rev)",
    "",
    "VIP từ sản phẩm phụ → revenue pillar",
    "Impact: +200-270M/tháng = +15-18% total rev",
], Inches(8.3), Inches(2.8), Inches(4.2), Inches(1.2),
              font_size=12, color=TEXT_WHITE)

# Key assumptions
add_text(slide, "Key Assumptions & Risks:", Inches(0.6), Inches(4.5), Inches(12), Inches(0.3),
         font_size=14, bold=True, color=BG_DARK)

add_table(slide, [
    ["Repeat rate 17% → 50%", "Level System + loss aversion", "Cao",
     "Core mechanic, đã proven ở nhiều game"],
    ["ARPT 1 → 3", "Streak + countdown + benefit tăng", "Trung bình",
     "Cần VIP value đủ mạnh ở Lv2-3"],
    ["PU 487 → 1,000", "Offer lần đầu + lapsed comeback", "Trung bình",
     "Phụ thuộc UA + conversion rate"],
    ["Bonus %G migration", "Remove shop → VIP exclusive", "Cao",
     "Existing payer buộc phải chuyển sang VIP"],
    ["Upgrade R via VIP", "Mở thêm đường cho mid user", "Trung bình",
     "Chỉ impact user có thẻ R chưa CLB B"],
], ["Assumption", "Mechanic", "Confidence", "Note"],
    Inches(0.3), Inches(5.0), Inches(12.7),
    col_widths=[2.5, 2.8, 1.2, 4.0], font_size=10)

# ============================================================
# SLIDE 9: Next Steps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = BG_DARK

add_text(slide, "Next Steps", Inches(1), Inches(1.5), Inches(11), Inches(0.8),
         font_size=36, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

steps = [
    "1. Chốt benefit table chính thức (review DKXX balance, Upgrade R rate)",
    "2. Estimate cost: bonus %G cho away bao nhiêu GEM/tháng?",
    "3. Design UI/UX: countdown timer, level progression screen, noti flow",
    "4. A/B test plan: test trên segment nhỏ trước khi roll out",
    "5. Monitor KPIs: repeat rate, ARPT, PU new, churn sau reset",
]
add_multiline(slide, steps, Inches(2), Inches(2.8), Inches(9), Inches(3.5),
              font_size=18, color=TEXT_WHITE, line_spacing=1.8)

add_shape(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), ACCENT)

# ============================================================
# Save
# ============================================================
out = "data/VIP_Solution_Concept_v1.pptx"
prs.save(out)
print(f"Saved: {out}")
